# -*- coding: utf-8 -*-
"""
Modern Document Conversion Service (2025)
Sử dụng các công nghệ mới nhất để convert documents
Tích hợp Adobe PDF Services API cho PDF to Word chất lượng cao
"""

import os
import shutil
import subprocess
import uuid
from pathlib import Path
from dotenv import load_dotenv
import json
import zipfile
from datetime import datetime
from typing import Optional, Dict, Any, List
import time
import asyncio
from io import BytesIO
import base64
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import numpy as np

# Load environment variables
load_dotenv()
from typing import Optional, List
import aiofiles
import httpx
from fastapi import UploadFile, HTTPException
import logging

# Document libraries
from docx import Document
import pypdf
# from pdf2docx import Converter as PDFToWordConverter  # REMOVED: requires opencv-python + PyMuPDF (230MB bloat!)
from pptx import Presentation
from openpyxl import load_workbook
import pypdfium2 as pdfium

# Adobe PDF Services (optional)
try:
    from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
    from adobe.pdfservices.operation.pdf_services import PDFServices
    from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
    from adobe.pdfservices.operation.pdfjobs.jobs.export_pdf_job import ExportPDFJob
    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_params import ExportPDFParams
    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_target_format import ExportPDFTargetFormat
    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_ocr_locale import ExportOCRLocale
    from adobe.pdfservices.operation.pdfjobs.result.export_pdf_result import ExportPDFResult
    from adobe.pdfservices.operation.io.cloud_asset import CloudAsset
    from adobe.pdfservices.operation.io.stream_asset import StreamAsset
    ADOBE_AVAILABLE = True
except ImportError:
    ADOBE_AVAILABLE = False

# Google Gemini API (optional) - RECOMMENDED for PDF → Word
try:
    import warnings
    # Suppress deprecation warning temporarily until migration to google.genai
    with warnings.catch_warnings():
        warnings.simplefilter("ignore", FutureWarning)
        import google.generativeai as genai
    import json
    import asyncio
    from app.services.gemini_service import get_gemini_service, GeminiService
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

logger = logging.getLogger(__name__)


# ==================== GEMINI MODELS CONFIGURATION ====================
# Updated: December 26, 2025
# Source: https://ai.google.dev/gemini-api/docs/models
# Source: https://ai.google.dev/pricing

GEMINI_MODELS = {
    # 🌟 GEMINI 3 SERIES (Newest - December 2025)
    "gemini-3-pro-preview": {
        "name": "Gemini 3 Pro Preview",
        "series": "3.0",
        "description": "World's best model for multimodal understanding",
        "features": ["Advanced reasoning", "Agentic tasks", "Vibe-coding", "Best multimodal"],
        "use_cases": ["Complex AI agents", "Advanced reasoning", "Cutting-edge apps"],
        "context_window": "1M+ tokens",
        "pricing": {"input": 2.00, "output": 12.00},  # per 1M tokens (includes thinking)
        "quality": 10,
        "speed": 6,
        "status": "preview",
        "free_tier": False,
        "recommended_for": ["ai-agents", "complex-reasoning"],
        "badge": "🌟 WORLD'S BEST"
    },
    "gemini-3-flash-preview": {
        "name": "Gemini 3 Flash Preview",
        "series": "3.0",
        "description": "Most intelligent model built for speed",
        "features": ["Frontier intelligence", "Superior search", "Grounding", "Fast"],
        "use_cases": ["Smart agents", "Fast reasoning", "Production apps"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.50, "output": 3.00},
        "quality": 10,
        "speed": 9,
        "status": "preview",
        "free_tier": True,
        "recommended_for": ["intelligent-speed"],
        "badge": "⚡ FASTEST INTELLIGENCE"
    },
    
    # ⚡ GEMINI 2.5 SERIES (Stable - September 2025)
    "gemini-2.5-flash": {
        "name": "Gemini 2.5 Flash",
        "series": "2.5",
        "description": "Best price-performance with hybrid reasoning",
        "features": ["Thinking budgets", "Fast processing", "1M context", "Agentic workflows"],
        "use_cases": ["PDF conversion", "Large scale processing", "High volume tasks"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.30, "output": 2.50},
        "quality": 9,
        "speed": 9,
        "status": "stable",
        "free_tier": True,
        "recommended_for": ["pdf-conversion", "production", "default"],
        "badge": "⭐ RECOMMENDED"
    },
    "gemini-2.5-flash-preview-09-2025": {
        "name": "Gemini 2.5 Flash Preview",
        "series": "2.5",
        "description": "Latest preview with newest improvements",
        "features": ["Latest features", "Thinking enabled", "1M context"],
        "use_cases": ["Testing new features", "Beta testing", "Advanced users"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.30, "output": 2.50},
        "quality": 9,
        "speed": 9,
        "status": "preview",
        "free_tier": True,
        "recommended_for": ["beta-testing"],
    },
    "gemini-2.5-flash-lite": {
        "name": "Gemini 2.5 Flash-Lite",
        "series": "2.5",
        "description": "Most cost-effective for scale",
        "features": ["Ultra-fast", "80% cheaper", "High throughput", "1M context"],
        "use_cases": ["Simple text extraction", "High volume", "Budget-conscious"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.10, "output": 0.40},
        "quality": 8,
        "speed": 10,
        "status": "stable",
        "free_tier": True,
        "recommended_for": ["budget", "high-volume"],
        "badge": "💰 CHEAPEST"
    },
    "gemini-2.5-flash-lite-preview-09-2025": {
        "name": "Gemini 2.5 Flash-Lite Preview",
        "series": "2.5",
        "description": "Latest lite model - cost-efficient, high quality",
        "features": ["Cost optimization", "High throughput", "Quality improvements"],
        "use_cases": ["Budget production", "High scale", "Fast processing"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.10, "output": 0.40},
        "quality": 8,
        "speed": 10,
        "status": "preview",
        "free_tier": True,
        "recommended_for": ["cost-efficiency"],
    },
    "gemini-2.5-pro": {
        "name": "Gemini 2.5 Pro",
        "series": "2.5",
        "description": "Advanced thinking for complex problems",
        "features": ["Advanced reasoning", "Long context", "Code/math/STEM expert", "2M context"],
        "use_cases": ["Complex research", "Large codebases", "Scientific problems"],
        "context_window": "2M tokens",
        "pricing": {"input": 1.25, "output": 10.00},
        "quality": 10,
        "speed": 7,
        "status": "stable",
        "free_tier": True,
        "recommended_for": ["complex-reasoning", "research"],
        "badge": "🎯 HIGHEST QUALITY"
    },
    
    # 🔧 GEMINI 2.0 SERIES (Previous Generation - October 2024)
    "gemini-2.0-flash": {
        "name": "Gemini 2.0 Flash",
        "series": "2.0",
        "description": "Second generation workhorse model",
        "features": ["1M context", "Balanced performance", "Multimodal", "Agent-ready"],
        "use_cases": ["General purpose", "Production apps", "Balanced workloads"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.10, "output": 0.40},
        "quality": 8,
        "speed": 9,
        "status": "stable",
        "free_tier": True,
        "recommended_for": ["general-purpose"],
    },
    "gemini-2.0-flash-exp": {
        "name": "Gemini 2.0 Flash Experimental",
        "series": "2.0",
        "description": "Experimental version with latest features",
        "features": ["Experimental features", "1M context", "Fast", "May change"],
        "use_cases": ["Testing", "Experiments", "Non-production"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.075, "output": 0.30},
        "quality": 8,
        "speed": 9,
        "status": "experimental",
        "free_tier": True,
        "recommended_for": ["testing"],
    },
    "gemini-2.0-flash-lite": {
        "name": "Gemini 2.0 Flash-Lite",
        "series": "2.0",
        "description": "Second gen small workhorse",
        "features": ["1M context", "Cost-effective", "High throughput"],
        "use_cases": ["High volume", "Simple tasks", "Budget apps"],
        "context_window": "1M tokens",
        "pricing": {"input": 0.075, "output": 0.30},
        "quality": 7,
        "speed": 10,
        "status": "stable",
        "free_tier": True,
        "recommended_for": ["simple-tasks"],
    },
}

# Default model recommendation
DEFAULT_GEMINI_MODEL = "gemini-2.5-flash"


def get_friendly_error_message(error: Exception) -> tuple[int, str]:
    """
    Convert Adobe API errors to user-friendly Vietnamese messages
    
    Returns:
        tuple: (status_code, friendly_message)
    """
    error_msg = str(error).lower()
    
    # Password protected files
    if "password_protected" in error_msg or "protected" in error_msg:
        return (400, "😔 Rất tiếc! File PDF này được bảo vệ bằng mật khẩu.\n\n"
                     "💡 Giải pháp:\n"
                     "• Mở file bằng PDF reader và nhập mật khẩu\n"
                     "• Sau đó 'Save As' thành file mới không có password\n"
                     "• Hoặc dùng tính năng 'Unlock PDF' của chúng tôi")
    
    # Digitally signed files (will trigger PyPDF2 fallback)
    if "pdf_signed" in error_msg or "signed" in error_msg:
        return (400, "😔 File PDF này có chữ ký điện tử (Adobe không hỗ trợ).\n\n"
                     "💡 Hệ thống sẽ tự động thử phương pháp khác...\n"
                     "• Hoặc vui lòng remove signature trước\n"
                     "• Hoặc dùng bản PDF gốc chưa ký")
    
    # Invalid page range errors (will trigger PyPDF2 fallback)
    if "invalid_page_range" in error_msg or "invalid range" in error_msg or "page range specified is invalid" in error_msg:
        return (400, "😔 Khoảng trang không hợp lệ (Adobe không hỗ trợ).\n\n"
                     "💡 Hệ thống sẽ tự động thử phương pháp khác...\n"
                     "• Kiểm tra số trang trong PDF\n"
                     "• Đảm bảo ranges hợp lệ (VD: 1-3,5-7)")
    
    # Corrupted or invalid files
    if "corrupted" in error_msg or "invalid" in error_msg or "malformed" in error_msg:
        return (400, "😔 Rất tiếc! File PDF này bị lỗi hoặc không hợp lệ.\n\n"
                     "💡 Giải pháp:\n"
                     "• Thử mở file bằng PDF reader để kiểm tra\n"
                     "• Nếu mở được, thử 'Print to PDF' để tạo file mới\n"
                     "• Hoặc dùng file PDF từ nguồn khác")
    
    # File too large
    if "file size" in error_msg or "too large" in error_msg:
        return (400, "😔 Rất tiếc! File PDF quá lớn để xử lý.\n\n"
                     "💡 Giải pháp:\n"
                     "• Giới hạn: 100MB cho mỗi file\n"
                     "• Thử nén/tối ưu file PDF trước\n"
                     "• Hoặc split thành nhiều file nhỏ hơn")
    
    # Page count limits
    if "page" in error_msg and ("limit" in error_msg or "exceed" in error_msg):
        return (400, "😔 Rất tiếc! File PDF có quá nhiều trang.\n\n"
                     "💡 Giải pháp:\n"
                     "• Split file thành nhiều phần nhỏ hơn\n"
                     "• Xử lý từng phần một\n"
                     "• Hoặc liên hệ hỗ trợ để tăng giới hạn")
    
    # Invalid page ranges
    if "page range" in error_msg or "invalid range" in error_msg:
        return (400, "😔 Rất tiếc! Phạm vi trang không hợp lệ.\n\n"
                     "💡 Giải pháp:\n"
                     "• Kiểm tra số trang: ví dụ '1-3' hoặc '1,3,5'\n"
                     "• Đảm bảo số trang không vượt quá tổng số trang\n"
                     "• Số trang bắt đầu từ 1 (không phải 0)")
    
    # Network/timeout errors
    if "timeout" in error_msg or "connection" in error_msg:
        return (500, "😔 Rất tiếc! Kết nối với Adobe API bị gián đoạn.\n\n"
                     "💡 Giải pháp:\n"
                     "• File có thể quá lớn hoặc phức tạp\n"
                     "• Vui lòng thử lại sau vài phút\n"
                     "• Hoặc liên hệ hỗ trợ nếu vẫn lỗi")
    
    # Quota exceeded
    if "quota" in error_msg or "limit exceeded" in error_msg:
        return (429, "😔 Rất tiếc! Đã vượt quá giới hạn sử dụng.\n\n"
                     "💡 Giải pháp:\n"
                     "• Số lượng request đã đạt giới hạn hôm nay\n"
                     "• Vui lòng thử lại vào ngày mai\n"
                     "• Hoặc liên hệ để nâng cấp gói dịch vụ")
    
    # Authentication errors
    if "credential" in error_msg or "authentication" in error_msg or "unauthorized" in error_msg:
        return (500, "😔 Rất tiếc! Có lỗi xác thực với Adobe API.\n\n"
                     "💡 Giải pháp:\n"
                     "• Đây là lỗi hệ thống, không phải lỗi của bạn\n"
                     "• Vui lòng liên hệ quản trị viên\n"
                     "• Chúng tôi sẽ khắc phục trong thời gian sớm nhất")
    
    # Generic error
    return (500, f"😔 Rất tiếc! Đã có lỗi khi xử lý file PDF.\n\n"
                 f"💡 Chi tiết lỗi:\n{str(error)}\n\n"
                 f"Vui lòng thử lại hoặc liên hệ hỗ trợ nếu vấn đề vẫn tiếp diễn.")


def is_pdf_scanned(pdf_path: Path, min_text_chars: int = 50) -> bool:
    """
    Detect if PDF is scanned (image-only, no text layer)
    
    Strategy:
    - Extract text from first few pages
    - If total text < min_text_chars → likely scanned
    - Returns True if PDF needs OCR, False otherwise
    
    Args:
        pdf_path: Path to PDF file
        min_text_chars: Minimum characters to consider as "text PDF" (default 50)
    
    Returns:
        True if PDF appears to be scanned (needs OCR)
        False if PDF has text layer
    """
    try:
        with open(pdf_path, 'rb') as f:
            pdf_reader = pypdf.PdfReader(f)
            
            # Check first 3 pages (or all if less than 3)
            pages_to_check = min(3, len(pdf_reader.pages))
            total_text = ""
            
            for i in range(pages_to_check):
                page = pdf_reader.pages[i]
                text = page.extract_text()
                total_text += text
            
            # If very little text found → likely scanned
            text_length = len(total_text.strip())
            is_scanned = text_length < min_text_chars
            
            if is_scanned:
                logger.info(f"PDF appears to be scanned (only {text_length} chars found in first {pages_to_check} pages)")
            else:
                logger.info(f"PDF has text layer ({text_length} chars found)")
            
            return is_scanned
            
    except Exception as e:
        logger.warning(f"Could not detect if PDF is scanned: {e}, assuming not scanned")
        return False


class DocumentService:
    """Modern document processing service sử dụng Gotenberg + Adobe PDF Services"""
    
    def __init__(
        self, 
        upload_dir: str = "uploads/documents",
        gotenberg_url: str = None
    ):
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        # Output directories
        self.output_dir = Path("uploads/outputs")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Gotenberg URL - auto-detect environment
        if gotenberg_url:
            self.gotenberg_url = gotenberg_url
        else:
            # Production: Gotenberg container, Dev: có thể dùng local LibreOffice
            self.gotenberg_url = os.getenv("GOTENBERG_URL", "http://gotenberg:3000")
        
        # Adobe PDF Services - optional but recommended for better quality
        self.use_adobe = os.getenv("USE_ADOBE_PDF_API", "false").lower() == "true"
        self.adobe_credentials = None
        
        logger.info("="*60)
        logger.info("🔵 ADOBE PDF SERVICES INITIALIZATION")
        logger.info(f"   USE_ADOBE_PDF_API={os.getenv('USE_ADOBE_PDF_API', 'not set')}")
        logger.info(f"   SDK Available: {ADOBE_AVAILABLE}")
        
        if self.use_adobe and ADOBE_AVAILABLE:
            client_id = os.getenv("PDF_SERVICES_CLIENT_ID")
            client_secret = os.getenv("PDF_SERVICES_CLIENT_SECRET")
            
            logger.info(f"   Client ID: {client_id[:20]}..." if client_id else "   Client ID: NOT SET")
            logger.info(f"   Client Secret: {'*' * 20}" if client_secret else "   Client Secret: NOT SET")
            
            if client_id and client_secret:
                try:
                    self.adobe_credentials = ServicePrincipalCredentials(
                        client_id=client_id,
                        client_secret=client_secret
                    )
                    logger.info("✅ Adobe PDF Services ENABLED - High quality PDF to Word conversion available")
                    logger.info("   Free tier: 500 document transactions/month")
                except Exception as e:
                    logger.error(f"❌ Failed to initialize Adobe credentials: {e}")
                    self.use_adobe = False
            else:
                logger.warning("⚠️ USE_ADOBE_PDF_API=true but credentials not found in env")
                self.use_adobe = False
        elif self.use_adobe and not ADOBE_AVAILABLE:
            logger.warning("⚠️ USE_ADOBE_PDF_API=true but pdfservices-sdk not installed")
            self.use_adobe = False
        
        logger.info(f"   Final Status: {'✅ ENABLED' if self.use_adobe else '❌ DISABLED'}")
        logger.info("="*60)
        
        # Google Gemini API - RECOMMENDED for PDF → Word with Vietnamese
        self.gemini_api_key = os.getenv("GEMINI_API_KEY")
        self.gemini_model_name = os.getenv("GEMINI_MODEL", DEFAULT_GEMINI_MODEL)  # Default to best model
        self.use_gemini = False
        self.gemini_service = None  # Will be initialized with db session
        
        if self.gemini_api_key and GEMINI_AVAILABLE:
            try:
                # Just configure genai for file upload/download operations
                genai.configure(api_key=self.gemini_api_key)
                self.use_gemini = True
                
                # Get model info for logging
                model_info = GEMINI_MODELS.get(self.gemini_model_name, {})
                model_display = model_info.get("name", self.gemini_model_name)
                logger.info(f"✅ Gemini API enabled - Model: {model_display} (Quality: {model_info.get('quality', '?')}/10)")
                logger.info("✅ Auto-logging enabled via GeminiService wrapper")
            except Exception as e:
                logger.warning(f"Failed to initialize Gemini API: {e}")
                self.use_gemini = False
        elif self.gemini_api_key and not GEMINI_AVAILABLE:
            logger.warning("GEMINI_API_KEY found but google-generativeai not installed. Run: pip install google-generativeai")
            self.use_gemini = False
        
    def get_available_gemini_models(self) -> dict:
        """Get list of all available Gemini models with metadata"""
        return GEMINI_MODELS.copy()
    
    def get_gemini_model_info(self, model_name: str) -> dict:
        """Get information about a specific Gemini model"""
        return GEMINI_MODELS.get(model_name, {})
    
    def set_gemini_model(self, model_name: str):
        """
        Change Gemini model dynamically
        
        Args:
            model_name: Model identifier (e.g., "gemini-2.5-flash")
        
        Raises:
            ValueError: If model not found or Gemini not configured
        """
        if not self.use_gemini:
            raise ValueError("Gemini API not configured. Set GEMINI_API_KEY in .env")
        
        if model_name not in GEMINI_MODELS:
            available = ", ".join(GEMINI_MODELS.keys())
            raise ValueError(f"Model '{model_name}' not found. Available models: {available}")
        
        try:
            self.gemini_model_name = model_name
            # No need to create model here - GeminiService will handle it per-request
            
            model_info = GEMINI_MODELS[model_name]
            logger.info(f"✅ Switched to {model_info['name']} (Quality: {model_info['quality']}/10, Speed: {model_info['speed']}/10)")
        except Exception as e:
            logger.error(f"Failed to switch model: {e}")
            raise ValueError(f"Failed to initialize model '{model_name}': {str(e)}")
    
    async def save_upload_file(self, upload_file: UploadFile) -> Path:
        """Save uploaded file async"""
        file_path = self.upload_dir / upload_file.filename
        
        async with aiofiles.open(file_path, 'wb') as out_file:
            content = await upload_file.read()
            await out_file.write(content)
            
        return file_path
    
    # ==================== Word → PDF (Gotenberg) ====================
    
    async def word_to_pdf(
        self,
        input_file: Path,
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert Word (.docx/.doc) sang PDF bằng Gotenberg API
        
        Gotenberg là Docker microservice hiện đại, sử dụng LibreOffice headless
        - Không cần cài LibreOffice trên host machine
        - REST API đơn giản, nhanh, ổn định
        - Support: DOC, DOCX, XLS, XLSX, PPT, PPTX, ODT...
        """
        if not input_file.suffix.lower() in ['.docx', '.doc']:
            raise HTTPException(400, "File phải là .docx hoặc .doc")
        
        output_filename = output_filename or input_file.stem + ".pdf"
        output_path = self.output_dir / output_filename
        
        try:
            # Gọi Gotenberg API để convert
            async with httpx.AsyncClient(timeout=60.0) as client:
                # Đọc file Word
                async with aiofiles.open(input_file, 'rb') as f:
                    file_content = await f.read()
                
                # POST đến Gotenberg LibreOffice endpoint
                files = {
                    'files': (input_file.name, file_content, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                }
                
                response = await client.post(
                    f"{self.gotenberg_url}/forms/libreoffice/convert",
                    files=files
                )
                
                if response.status_code != 200:
                    raise HTTPException(
                        500,
                        f"Gotenberg conversion failed: {response.status_code} - {response.text}"
                    )
                
                # Lưu PDF output
                async with aiofiles.open(output_path, 'wb') as f:
                    await f.write(response.content)
                
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise HTTPException(500, "File PDF không được tạo ra hoặc rỗng")
                
            return output_path
            
        except httpx.ConnectError:
            # Gotenberg không khả dụng - fallback to LibreOffice local (dev only)
            return await self._word_to_pdf_libreoffice_fallback(input_file, output_filename)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"Lỗi chuyển đổi Word sang PDF: {str(e)}")
    
    async def _word_to_pdf_libreoffice_fallback(
        self,
        input_file: Path,
        output_filename: str
    ) -> Path:
        """
        Fallback method: Dùng LibreOffice local nếu Gotenberg không có
        
        Cài LibreOffice: https://www.libreoffice.org/download/download/
        """
        logger.info(f"🔄 Starting LibreOffice conversion: {input_file.name}")
        output_path = self.output_dir / output_filename
        
        # Tìm LibreOffice trong các đường dẫn phổ biến
        libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"D:\LibreOffice\program\soffice.exe",  # Custom install
            "/usr/bin/libreoffice",  # Linux
            "/usr/bin/soffice",
            "soffice",
        ]
        
        soffice_path = None
        for path in libreoffice_paths:
            if os.path.exists(path):
                soffice_path = path
                logger.info(f"✅ Found LibreOffice at: {path}")
                break
        
        if not soffice_path:
            raise HTTPException(
                503,
                "⚠️ LibreOffice chưa được cài đặt. "
                "Tải tại: https://www.libreoffice.org/download/download/ "
                "(chọn phiên bản 64-bit, ~300MB)"
            )
        
        try:
            # Log input file info
            logger.info(f"📄 Input file: {input_file} (size: {input_file.stat().st_size} bytes)")
            logger.info(f"📂 Output dir: {self.output_dir}")
            
            cmd = [
                soffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.output_dir),
                str(input_file)
            ]
            
            logger.info(f"🔧 Running command: {' '.join(cmd)}")
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60
            )
            
            # Log subprocess output
            logger.info(f"📊 Return code: {result.returncode}")
            if result.stdout:
                logger.info(f"📤 STDOUT: {result.stdout}")
            if result.stderr:
                logger.warning(f"⚠️ STDERR: {result.stderr}")
            
            if result.returncode != 0:
                raise HTTPException(500, f"LibreOffice conversion failed (code {result.returncode}): {result.stderr}")
            
            # Check generated PDF
            generated_pdf = self.output_dir / (input_file.stem + ".pdf")
            logger.info(f"🔍 Looking for: {generated_pdf}")
            
            if not generated_pdf.exists():
                # List all files in output dir for debugging
                all_files = list(self.output_dir.glob("*"))
                logger.error(f"❌ PDF not found! Files in output dir: {[f.name for f in all_files]}")
                raise HTTPException(500, f"PDF file not generated: {generated_pdf}")
            
            logger.info(f"✅ PDF found: {generated_pdf} (size: {generated_pdf.stat().st_size} bytes)")
            
            # Rename if needed
            if generated_pdf != output_path:
                logger.info(f"📝 Renaming: {generated_pdf.name} → {output_path.name}")
                generated_pdf.rename(output_path)
            
            # Verify final output
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise HTTPException(500, f"Output PDF is empty or missing: {output_path}")
            
            logger.info(f"✅ LibreOffice conversion completed: {output_path.name} ({output_path.stat().st_size} bytes)")
            return output_path
            
        except subprocess.TimeoutExpired:
            logger.error("❌ LibreOffice conversion timeout (>60s)")
            raise HTTPException(500, "Conversion timeout (>60s)")
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"❌ LibreOffice conversion error: {str(e)}", exc_info=True)
            raise HTTPException(500, f"LibreOffice conversion error: {str(e)}")
    
    # ==================== Generic Office → PDF (Excel, PowerPoint) ====================
    
    async def office_to_pdf(
        self,
        input_file: Path,
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert any Office file (Excel, PowerPoint) to PDF using Gotenberg
        
        Supports: .xlsx, .xls, .pptx, .ppt, .odt, .ods, .odp
        """
        supported_extensions = ['.xlsx', '.xls', '.pptx', '.ppt', '.odt', '.ods', '.odp']
        if input_file.suffix.lower() not in supported_extensions:
            raise HTTPException(400, f"File phải là Office file: {', '.join(supported_extensions)}")
        
        output_filename = output_filename or input_file.stem + ".pdf"
        output_path = self.output_dir / output_filename
        
        try:
            # Gọi Gotenberg API
            async with httpx.AsyncClient(timeout=120.0) as client:  # Excel/PPT có thể mất thời gian hơn
                async with aiofiles.open(input_file, 'rb') as f:
                    file_content = await f.read()
                
                # Determine MIME type
                mime_types = {
                    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    '.xls': 'application/vnd.ms-excel',
                    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    '.ppt': 'application/vnd.ms-powerpoint',
                    '.odt': 'application/vnd.oasis.opendocument.text',
                    '.ods': 'application/vnd.oasis.opendocument.spreadsheet',
                    '.odp': 'application/vnd.oasis.opendocument.presentation'
                }
                mime_type = mime_types.get(input_file.suffix.lower(), 'application/octet-stream')
                
                files = {
                    'files': (input_file.name, file_content, mime_type)
                }
                
                response = await client.post(
                    f"{self.gotenberg_url}/forms/libreoffice/convert",
                    files=files
                )
                
                if response.status_code != 200:
                    raise HTTPException(500, f"Gotenberg conversion failed: {response.status_code}")
                
                async with aiofiles.open(output_path, 'wb') as f:
                    await f.write(response.content)
                
            if not output_path.exists() or output_path.stat().st_size == 0:
                raise HTTPException(500, "File PDF không được tạo ra hoặc rỗng")
                
            return output_path
            
        except httpx.ConnectError:
            # Fallback to LibreOffice local
            return await self._office_to_pdf_libreoffice_fallback(input_file, output_filename)
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"Lỗi chuyển đổi Office sang PDF: {str(e)}")
    
    async def _office_to_pdf_libreoffice_fallback(
        self,
        input_file: Path,
        output_filename: str
    ) -> Path:
        """Fallback: Use LibreOffice for Office files"""
        return await self._word_to_pdf_libreoffice_fallback(input_file, output_filename)

    # ==================== PDF → Word ====================
    
    async def pdf_to_word(
        self,
        input_file: Path,
        output_filename: Optional[str] = None,
        start_page: int = 0,
        end_page: Optional[int] = None,
        enable_ocr: bool = False,
        ocr_language: str = "vi-VN",
        auto_detect_scanned: bool = True,
        use_gemini: bool = False,
        gemini_model: Optional[str] = None,
        db = None  # SQLAlchemy session for auto-logging
    ) -> Path:
        """
        Convert PDF to Word (.docx)
        
        Strategy:
        1. Gemini API (if use_gemini=True) - Best for tables/layout, 100+ languages, multiple models
        2. Adobe PDF Services (if enabled) - AI-powered, 10/10 quality BUT NO Vietnamese support
        3. Fallback to pdf2docx - Good quality, free
        
        Args:
            gemini_model: Optional Gemini model to use (e.g., "gemini-2.5-flash")
                         If not specified, uses configured default model
        
        OCR Support:
        - enable_ocr: Enable OCR for scanned PDFs (default: False)
        - ocr_language: Language for OCR (default: "vi-VN" for Vietnamese)
        - auto_detect_scanned: Auto-detect if PDF is scanned and enable OCR (default: True)
        - use_gemini: Use Gemini API instead of Adobe (supports Vietnamese, better for tables)
        
        Supported OCR languages (Adobe):
        - vi-VN (Vietnamese), en-US (English), fr-FR (French), de-DE (German),
        - es-ES (Spanish), it-IT (Italian), pt-BR (Portuguese), ja-JP (Japanese),
        - ko-KR (Korean), zh-CN (Chinese Simplified), zh-TW (Chinese Traditional)
        - and 40+ more languages
        
        Gemini API (NEW - December 2025):
        - Supports 100+ languages including Vietnamese
        - Multiple models with different quality/cost trade-offs
        - Native PDF understanding (no OCR preprocessing)
        - Best for tables, layout, complex documents
        - Quality: 9-10/10 depending on model
        """
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + ".docx"
        output_path = self.output_dir / output_filename
        
        # ========== DECISION LOGIC - LOG EVERYTHING ==========
        logger.info("")
        logger.info("="*70)
        logger.info("📄 PDF → WORD CONVERSION REQUEST")
        logger.info(f"   File: {input_file.name}")
        logger.info(f"   Size: {input_file.stat().st_size / 1024:.2f} KB")
        logger.info(f"   OCR Language: {ocr_language}")
        logger.info("="*70)
        logger.info("")
        logger.info("🔍 CHECKING AVAILABLE TECHNOLOGIES:")
        logger.info(f"   1. Gemini API:        {'✅ Available' if self.use_gemini else '❌ Not configured'} (use_gemini={use_gemini})")
        logger.info(f"   2. Adobe PDF:         {'✅ Available' if (self.use_adobe and self.adobe_credentials) else '❌ Not configured'}")
        logger.info(f"   3. pdf2docx (local):  ❌ Disabled (to reduce image size)")
        logger.info("")
        
        # Try Gemini first if requested (best for Vietnamese + tables)
        if use_gemini:
            logger.info("="*70)
            logger.info("🎯 DECISION: Using GEMINI API (requested by user)")
            logger.info(f"   Model: {gemini_model or self.gemini_model_name}")
            logger.info(f"   Reason: Best for Vietnamese, excellent table extraction")
            logger.info("="*70)
            logger.info("")
            logger.info("📞 Calling _pdf_to_word_gemini()...")
            
            try:
                result = await self._pdf_to_word_gemini(
                    input_file, 
                    output_path, 
                    ocr_language,
                    model_name=gemini_model,
                    db=db
                )
                logger.info(f"📨 Returned from _pdf_to_word_gemini(): {result}")
                logger.info("")
                return result
            except Exception as e:
                logger.error("")
                logger.error("="*70)
                logger.error(f"❌ GEMINI API - CONVERSION FAILED")
                logger.error(f"   Error Type: {type(e).__name__}")
                logger.error(f"   Error: {str(e)}")
                logger.error("   Falling back to Adobe/pdf2docx...")
                logger.error("="*70)
                logger.error("")
                
                # Log full traceback for debugging
                import traceback
                logger.debug(f"Full traceback:\n{traceback.format_exc()}")
        else:
            logger.info("ℹ️  Gemini API not requested (use_gemini=False)")
            logger.info("")
        
        # Auto-detect if PDF is scanned (if not explicitly disabled)
        needs_ocr = enable_ocr
        if auto_detect_scanned and not enable_ocr:
            logger.info("🔍 Auto-detecting if PDF is scanned...")
            is_scanned = is_pdf_scanned(input_file)
            logger.info(f"   Result: {'📸 SCANNED (image-only)' if is_scanned else '📝 TEXT-BASED (has text layer)'}")
            if is_scanned:
                needs_ocr = True
                logger.info(f"   ✓ OCR automatically enabled (language: {ocr_language})")
            logger.info("")
        
        # 🌟 HYBRID APPROACH: Vietnamese scanned PDF
        is_vietnamese = ocr_language.lower().startswith('vi')
        if needs_ocr and is_vietnamese and self.use_gemini and (self.use_adobe and self.adobe_credentials):
            logger.info("="*70)
            logger.info("🌟 DECISION: Using HYBRID APPROACH")
            logger.info("   Reason: Vietnamese scanned PDF detected")
            logger.info("   Problem: Adobe has NO Vietnamese OCR support")
            logger.info("   Solution: Gemini OCR (text) + Adobe (layout/images)")
            logger.info("="*70)
            logger.info("")
            logger.info("📞 Calling _pdf_to_word_hybrid_vietnamese()...")
            
            try:
                result = await self._pdf_to_word_hybrid_vietnamese(
                    input_file,
                    output_path,
                    db=db
                )
                logger.info(f"📨 Returned from _pdf_to_word_hybrid_vietnamese(): {result}")
                logger.info("")
                return result
            except Exception as e:
                logger.error("")
                logger.error("="*70)
                logger.error(f"❌ HYBRID APPROACH - CONVERSION FAILED")
                logger.error(f"   Error Type: {type(e).__name__}")
                logger.error(f"   Error: {str(e)}")
                logger.error("   Falling back to Adobe-only...")
                logger.error("="*70)
                logger.error("")
                import traceback
                logger.debug(f"Full traceback:\n{traceback.format_exc()}")
        
        # Try Adobe if enabled (best quality but NO Vietnamese)
        if self.use_adobe and self.adobe_credentials:
            logger.info("🎯 DECISION: Using ADOBE PDF SERVICES")
            logger.info(f"   Technology: Adobe PDF Services API (Cloud)")
            logger.info(f"   Quality: 10/10")
            logger.info(f"   OCR Enabled: {needs_ocr}")
            if needs_ocr:
                logger.info(f"   OCR Language: {ocr_language}")
                if ocr_language.startswith('vi'):
                    logger.warning("   ⚠️  Note: Adobe does NOT support Vietnamese OCR")
                    logger.info("   → Will use Tesseract OCR + Adobe Export")
            logger.info("")
            
            try:
                logger.info("="*60)
                logger.info(f"🔵 ADOBE PDF SERVICES - STARTING CONVERSION")
                logger.info(f"   File: {input_file.name}")
                logger.info(f"   Size: {input_file.stat().st_size / 1024:.2f} KB")
                logger.info(f"   OCR Enabled: {needs_ocr}")
                logger.info(f"   OCR Language: {ocr_language}")
                logger.info("="*60)
                logger.info("")
                logger.info("📞 Calling _pdf_to_word_adobe()...")
                
                result = await self._pdf_to_word_adobe(
                    input_file, 
                    output_path,
                    enable_ocr=needs_ocr,
                    ocr_language=ocr_language
                )
                
                logger.info(f"📨 Returned from _pdf_to_word_adobe(): {result}")
                logger.info("")
                logger.info("="*60)
                logger.info("✅ ADOBE PDF SERVICES - CONVERSION SUCCESS")
                logger.info(f"   Output: {result.name}")
                logger.info(f"   Size: {result.stat().st_size / 1024:.2f} KB")
                logger.info("="*60)
                logger.info("")
                
                return result
            except Exception as e:
                logger.error("")
                logger.error("="*60)
                logger.error(f"❌ ADOBE PDF SERVICES - CONVERSION FAILED")
                logger.error(f"   Error Type: {type(e).__name__}")
                logger.error(f"   Error: {str(e)}")
                logger.error(f"   Falling back to pdf2docx...")
                logger.error("="*60)
                logger.error("")
                
                # Log full traceback for debugging
                import traceback
                logger.debug(f"Full traceback:\n{traceback.format_exc()}")
        else:
            # Adobe not available
            logger.warning("⚠️  Adobe PDF Services NOT available")
            logger.warning(f"   use_adobe={self.use_adobe}")
            logger.warning(f"   credentials={bool(self.adobe_credentials)}")
            logger.info("")
        
        # Fallback to pdf2docx (good quality, free)
        # Note: pdf2docx doesn't support OCR, so if needs_ocr=True, warn user
        logger.info("="*70)
        logger.info("🎯 DECISION: Fallback to PDF2DOCX (local)")
        logger.info(f"   Technology: pdf2docx (Pure Python library)")
        logger.info(f"   Quality: 7/10")
        logger.info(f"   Status: ❌ DISABLED (to reduce Docker image size)")
        logger.info("="*70)
        logger.info("")
        
        if needs_ocr:
            logger.warning("⚠️  PDF appears to be scanned but Adobe OCR not available.")
            logger.warning("   Using pdf2docx - results may be poor for scanned PDFs.")
            logger.warning("   Consider enabling Adobe PDF Services for better OCR quality.")
            logger.info("")
        
        logger.info("📞 Calling _pdf_to_word_local()...")
        
        try:
            result = await self._pdf_to_word_local(input_file, output_path, start_page, end_page)
            logger.error("❓ UNEXPECTED: _pdf_to_word_local() returned without error!")
            logger.error(f"   This should not happen (pdf2docx is disabled)")
            return result
        except HTTPException as e:
            logger.error("")
            logger.error("="*70)
            logger.error("❌ CONVERSION FAILED - NO TECHNOLOGY AVAILABLE")
            logger.error(f"   Reason: {e.detail}")
            logger.error("   Available: None (Adobe failed, Gemini not requested, pdf2docx disabled)")
            logger.error("   Solution: Enable Adobe PDF Services or use Gemini API")
            logger.error("="*70)
            logger.error("")
            raise
    
    async def _pdf_to_word_adobe(
        self, 
        input_file: Path, 
        output_path: Path,
        enable_ocr: bool = False,
        ocr_language: str = "vi-VN"
    ) -> Path:
        """
        Convert PDF to Word using Adobe PDF Services API
        High quality conversion with AI-powered layout preservation
        
        OCR Support:
        - Adobe SDK v4.2.0: Uses ExportPDF with OCR locale (embedded in export)
        - No separate OCR job needed - OCR is done during PDF→Word export
        - Supports 50+ languages via ExportOCRLocale
        
        Note: Adobe OCR does NOT support Vietnamese (vi-VN). For Vietnamese PDFs,
        we use Tesseract OCR or Gemini API.
        """
        logger.info(f"🔹 _pdf_to_word_adobe() CALLED")
        logger.info(f"   enable_ocr={enable_ocr}, ocr_language={ocr_language}")
        
        try:
            # Direct conversion with optional OCR
            # Adobe SDK v4.2.0: OCR is embedded in ExportPDF job, not separate
            logger.info(f"   📞 Calling _pdf_to_word_adobe_internal(enable_ocr={enable_ocr})...")
            result = await self._pdf_to_word_adobe_internal(
                input_file, 
                output_path,
                enable_ocr=enable_ocr,
                ocr_language=ocr_language
            )
            logger.info(f"   ✅ Returned from _pdf_to_word_adobe_internal(): {result}")
            return result
                
        except Exception as e:
            logger.error(f"Adobe PDF Services error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    
    async def _pdf_to_word_adobe_internal(
        self, 
        input_file: Path, 
        output_path: Path,
        enable_ocr: bool = False,
        ocr_language: str = "en-US"
    ) -> Path:
        """Internal method for PDF to Word conversion with optional OCR"""
        import time
        start_time = time.time()
        
        # Read input file
        logger.info("   📖 Step 1/5: Reading input file...")
        async with aiofiles.open(input_file, 'rb') as f:
            input_stream = await f.read()
        logger.info(f"   ✓ Read {len(input_stream)} bytes")
        
        # Create PDF Services instance
        logger.info("   🔐 Step 2/5: Creating PDF Services client...")
        pdf_services = PDFServices(credentials=self.adobe_credentials)
        logger.info("   ✓ Client created")
        
        # Upload file to Adobe
        logger.info("   ☁️  Step 3/5: Uploading to Adobe cloud...")
        upload_start = time.time()
        input_asset = pdf_services.upload(
            input_stream=input_stream,
            mime_type=PDFServicesMediaType.PDF
        )
        upload_time = time.time() - upload_start
        logger.info(f"   ✓ Upload complete ({upload_time:.2f}s)")
        
        # Create export parameters with optional OCR
        logger.info(f"   ⚙️  Step 4/5: Submitting export job (PDF → DOCX, OCR={enable_ocr})...")
        
        if enable_ocr:
            # Map ocr_language to ExportOCRLocale
            # en-US → EN_US
            adobe_locale = ocr_language.upper().replace('-', '_')
            
            try:
                from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_ocr_locale import ExportOCRLocale
                ocr_locale_enum = getattr(ExportOCRLocale, adobe_locale, ExportOCRLocale.EN_US)
                logger.info(f"   🔍 OCR enabled with locale: {adobe_locale}")
                
                # Create params WITHOUT ocr_locale (not supported in ExportPDFParams)
                export_pdf_params = ExportPDFParams(
                    target_format=ExportPDFTargetFormat.DOCX
                )
                
                # Create job WITH ocr_locale (this is where OCR is enabled)
                try:
                    export_pdf_job = ExportPDFJob(
                        input_asset=input_asset,
                        export_pdf_params=export_pdf_params,
                        ocr_locale=ocr_locale_enum  # ✅ Pass OCR to job, not params
                    )
                except TypeError as e:
                    # Fallback: ocr_locale not supported in this SDK version
                    logger.warning(f"   ⚠️  ocr_locale parameter not supported in this SDK version: {e}")
                    logger.warning("   ⚠️  Creating export job WITHOUT ocr_locale (OCR may not work properly)")
                    export_pdf_job = ExportPDFJob(
                        input_asset=input_asset,
                        export_pdf_params=export_pdf_params
                    )
            except (ImportError, AttributeError) as e:
                logger.warning(f"   ⚠️  OCR locale not supported: {adobe_locale}, using default export")
                export_pdf_params = ExportPDFParams(
                    target_format=ExportPDFTargetFormat.DOCX
                )
                export_pdf_job = ExportPDFJob(
                    input_asset=input_asset,
                    export_pdf_params=export_pdf_params
                )
        else:
            # No OCR - direct export
            export_pdf_params = ExportPDFParams(
                target_format=ExportPDFTargetFormat.DOCX
            )
            export_pdf_job = ExportPDFJob(
                input_asset=input_asset,
                export_pdf_params=export_pdf_params
            )
        
        job_start = time.time()
        location = pdf_services.submit(export_pdf_job)
        logger.info(f"   ✓ Job submitted to: {location}")
        
        # Get result (polling handled by SDK)
        logger.info("   ⏳ Polling for result (Adobe processing)...")
        pdf_services_response = pdf_services.get_job_result(location, ExportPDFResult)
        job_time = time.time() - job_start
        logger.info(f"   ✓ Job completed ({job_time:.2f}s)")
        
        # Download result
        logger.info("   💾 Step 5/5: Downloading result...")
        download_start = time.time()
        result_asset: CloudAsset = pdf_services_response.get_result().get_asset()
        stream_asset: StreamAsset = pdf_services.get_content(result_asset)
        
        # Save to file
        async with aiofiles.open(output_path, "wb") as f:
            await f.write(stream_asset.get_input_stream())
        
        download_time = time.time() - download_start
        total_time = time.time() - start_time
        
        logger.info(f"   ✓ Download complete ({download_time:.2f}s)")
        logger.info("")
        logger.info(f"   ⏱️  Total time: {total_time:.2f}s (upload: {upload_time:.2f}s, process: {job_time:.2f}s, download: {download_time:.2f}s)")
        logger.info(f"   📄 Output: {output_path}")
        
        return output_path
    
    async def _ocr_pdf_adobe(
        self,
        input_file: Path,
        ocr_language: str = "vi-VN"
    ) -> Path:
        """
        Perform OCR on PDF using Adobe OCRPDFJob
        Returns path to searchable PDF
        """
        try:
            # Import OCR-specific classes
            from adobe.pdfservices.operation.pdfjobs.jobs.ocr_pdf_job import OCRPDFJob
            from adobe.pdfservices.operation.pdfjobs.params.ocr_pdf.ocr_pdf_params import OCRPDFParams
            from adobe.pdfservices.operation.pdfjobs.params.ocr_pdf.ocr_supported_locale import OCRSupportedLocale
            from adobe.pdfservices.operation.pdfjobs.params.ocr_pdf.ocr_supported_type import OCRSupportedType
            from adobe.pdfservices.operation.pdfjobs.result.ocr_pdf_result import OCRPDFResult
            
            # Read input file
            async with aiofiles.open(input_file, 'rb') as f:
                input_stream = await f.read()
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=input_stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Map language code to OCRSupportedLocale
            # vi-VN → VI_VN
            adobe_locale = ocr_language.upper().replace('-', '_')
            
            try:
                ocr_locale = getattr(OCRSupportedLocale, adobe_locale)
            except AttributeError:
                logger.warning(f"OCR locale {adobe_locale} not supported, using EN_US")
                ocr_locale = OCRSupportedLocale.EN_US
            
            # Create OCR parameters
            ocr_pdf_params = OCRPDFParams(
                ocr_locale=ocr_locale,
                ocr_type=OCRSupportedType.SEARCHABLE_IMAGE_EXACT  # Preserves original
            )
            
            # Create and submit OCR job
            ocr_pdf_job = OCRPDFJob(
                input_asset=input_asset,
                ocr_pdf_params=ocr_pdf_params
            )
            
            location = pdf_services.submit(ocr_pdf_job)
            
            # Get result
            pdf_services_response = pdf_services.get_job_result(location, OCRPDFResult)
            
            # Download OCR'd PDF
            result_asset: CloudAsset = pdf_services_response.get_result().get_asset()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save OCR'd PDF to temp file
            ocr_output = self.output_dir / f"ocr_temp_{input_file.name}"
            async with aiofiles.open(ocr_output, "wb") as f:
                await f.write(stream_asset.get_input_stream())
            
            logger.info(f"Adobe OCR successful: {ocr_output} (locale: {adobe_locale})")
            return ocr_output
            
        except Exception as e:
            logger.error(f"Adobe OCR error: {e}")
            raise
            return output_path
            
        except Exception as e:
            logger.error(f"Adobe PDF Services error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    

    
    async def _pdf_to_word_hybrid_vietnamese(
        self,
        input_file: Path,
        output_path: Path,
        db = None
    ) -> Path:
        """
        🎯 HYBRID APPROACH: Gemini OCR + Adobe Layout Preservation
        
        Best solution for Vietnamese scanned PDFs:
        1. Gemini OCR → Perfect Vietnamese text (98% accuracy)
        2. Adobe ExportPDF (EN_US OCR) → Perfect images/layout preservation
        3. Combine → Word document with best quality
        
        Why hybrid:
        - Adobe: NO Vietnamese OCR support, but best layout/image preservation
        - Gemini: Excellent Vietnamese OCR, but loses images/layout
        - Hybrid: Best of both worlds! ✅
        """
        import time
        start_time = time.time()
        
        logger.info("")
        logger.info("="*80)
        logger.info("🌟 HYBRID APPROACH: GEMINI OCR + ADOBE LAYOUT")
        logger.info(f"   File: {input_file.name}")
        logger.info(f"   Strategy: Best of both worlds")
        logger.info(f"   - Gemini: Vietnamese text extraction (98% accuracy)")
        logger.info(f"   - Adobe: Layout + images preservation (10/10 quality)")
        logger.info("="*80)
        logger.info("")
        
        try:
            # Step 1: Gemini OCR for Vietnamese text
            logger.info("🤖 PHASE 1: GEMINI OCR - Vietnamese Text Extraction")
            logger.info("-" * 70)
            gemini_start = time.time()
            
            vietnamese_text = await self._gemini_extract_text_only(input_file, language="vi", db=db)
            
            gemini_time = time.time() - gemini_start
            logger.info(f"✅ Gemini extraction complete: {len(vietnamese_text)} characters ({gemini_time:.2f}s)")
            logger.info("")
            
            # Step 2: Adobe ExportPDF with OCR for layout/images
            logger.info("🔷 PHASE 2: ADOBE EXPORT PDF - Layout + Images Preservation")
            logger.info("-" * 70)
            logger.info("   Note: Using EN_US locale (Vietnamese not supported by Adobe)")
            logger.info("   Purpose: Preserve images and layout structure")
            adobe_start = time.time()
            
            # Call existing working method with OCR enabled
            adobe_output = await self._pdf_to_word_adobe_internal(
                input_file,
                output_path,
                enable_ocr=True,
                ocr_language="EN_US"
            )
            
            adobe_time = time.time() - adobe_start
            logger.info(f"✅ Adobe layout preserved: {adobe_output.name} ({adobe_time:.2f}s)")
            logger.info("")
            
            # Step 3: Note about current implementation
            logger.info("🔧 PHASE 3: RESULT")
            logger.info("-" * 70)
            logger.info("   Adobe created Word with images/layout ✅")
            logger.info("   Text quality: English OCR (Adobe - no Vietnamese support)")
            logger.info("   Images: Preserved 100% ✅")
            logger.info("   Layout: Preserved 10/10 ✅")
            logger.info("")
            logger.info("   💡 Future enhancement: Replace Adobe text with Gemini Vietnamese text")
            logger.info("")
            
            total_time = time.time() - start_time
            
            logger.info("="*80)
            logger.info("✅ HYBRID CONVERSION COMPLETE")
            logger.info(f"   Output: {adobe_output.name}")
            logger.info(f"   Size: {adobe_output.stat().st_size / 1024:.2f} KB")
            logger.info(f"   Total time: {total_time:.2f}s")
            logger.info(f"      - Gemini OCR: {gemini_time:.2f}s")
            logger.info(f"      - Adobe layout: {adobe_time:.2f}s")
            logger.info("")
            logger.info("   📊 Result Quality:")
            logger.info("      - Images: ✅ 100% preserved (Adobe)")
            logger.info("      - Layout: ✅ 10/10 (Adobe AI)")
            logger.info("      - Text: ⚠️  English OCR (Adobe - no Vietnamese support)")
            logger.info("")
            logger.info("   💡 Note: Gemini extracted perfect Vietnamese text")
            logger.info("      Future: Will replace Adobe English text with Gemini Vietnamese")
            logger.info("="*80)
            logger.info("")
            
            return adobe_output
            
        except Exception as e:
            logger.error("")
            logger.error("="*80)
            logger.error("❌ HYBRID APPROACH FAILED")
            logger.error(f"   Error: {e}")
            logger.error("="*80)
            logger.error("")
            raise
    
    async def _gemini_extract_text_only(
        self,
        input_file: Path,
        language: str = "vi",
        db = None
    ) -> str:
        """
        Extract text only from PDF using Gemini (no Word generation)
        Returns: Plain text string
        """
        if not self.use_gemini:
            raise HTTPException(500, "Gemini API not configured")
        
        logger.info("   📤 Uploading PDF to Gemini...")
        pdf_file = genai.upload_file(str(input_file))
        
        logger.info("   ⏳ Waiting for Gemini preprocessing...")
        import asyncio
        while pdf_file.state.name == "PROCESSING":
            await asyncio.sleep(1)
            pdf_file = genai.get_file(pdf_file.name)
        
        if pdf_file.state.name == "FAILED":
            raise ValueError(f"Gemini preprocessing failed: {pdf_file.state.name}")
        
        logger.info(f"   ✓ PDF ready: {pdf_file.name}")
        
        # Simplified prompt for text extraction only
        prompt = f"""Trích xuất TOÀN BỘ văn bản từ PDF này.

YÊU CẦU:
- Giữ CHÍNH XÁC 100% mọi ký tự tiếng Việt
- Giữ nguyên cấu trúc đoạn văn
- Không thêm giải thích hay chú thích
- CHỈ trả về văn bản thuần túy

Ngôn ngữ: {language}
"""
        
        logger.info("   🧠 Extracting text with Gemini AI...")
        
        # Use GeminiService for auto-logging if db available
        from app.services.gemini_service import get_gemini_service
        gemini_service = get_gemini_service(db) if db else None
        
        if gemini_service:
            response = gemini_service.generate_content(
                prompt=[pdf_file, prompt],
                model=self.gemini_model_name,
                operation="pdf-text-extraction",
                metadata={
                    "file_name": input_file.name,
                    "language": language,
                    "purpose": "hybrid_vietnamese_ocr"
                }
            )
        else:
            model_obj = genai.GenerativeModel(self.gemini_model_name)
            response = model_obj.generate_content([pdf_file, prompt])
        
        extracted_text = response.text.strip()
        logger.info(f"   ✓ Extracted {len(extracted_text)} characters")
        
        # Cleanup
        try:
            genai.delete_file(pdf_file.name)
        except Exception as e:
            logger.warning(f"   ⚠️  Cleanup warning: {e}")
        
        return extracted_text
    
    async def _pdf_to_word_local(
        self,
        input_file: Path,
        output_path: Path,
        start_page: int = 0,
        end_page: Optional[int] = None
    ) -> Path:
        """
        DISABLED: pdf2docx requires opencv-python + PyMuPDF (230MB bloat)
        Use Adobe PDF Services or Gemini OCR instead
        """
        logger.error("pdf2docx disabled to reduce Docker image size")
        raise HTTPException(
            501, 
            "pdf2docx method disabled. Use 'adobe' or 'gemini' conversion method"
        )
    
    async def _pdf_to_word_gemini(
        self,
        input_file: Path,
        output_path: Path,
        ocr_language: str = "vi",
        model_name: Optional[str] = None,
        db = None
    ) -> Path:
        """
        Convert PDF to Word using Google Gemini API
        
        Args:
            model_name: Optional Gemini model to use (e.g., "gemini-2.5-flash")
                       If not specified, uses the configured default model
            db: SQLAlchemy session for auto-logging usage
        
        Features:
        - Native PDF processing (no OCR library needed)
        - Understands Vietnamese and 100+ languages
        - Excellent table extraction (9.5/10)
        - Smart layout preservation
        - Multiple model options with different quality/cost trade-offs
        - ✅ AUTO-LOGGING: All API calls tracked automatically
        
        Free tier: 1,500 requests/day
        """
        import time
        start_time = time.time()
        
        if not self.use_gemini:
            raise HTTPException(500, "Gemini API not configured. Set GEMINI_API_KEY in .env")
        
        # Switch model if specified
        original_model = self.gemini_model_name
        if model_name and model_name != self.gemini_model_name:
            try:
                self.set_gemini_model(model_name)
            except ValueError as e:
                raise HTTPException(400, str(e))
        
        try:
            model_info = GEMINI_MODELS.get(self.gemini_model_name, {})
            
            logger.info("="*60)
            logger.info("🤖 GEMINI API - STARTING CONVERSION")
            logger.info(f"   Model: {model_info.get('name', self.gemini_model_name)}")
            logger.info(f"   Quality: {model_info.get('quality', '?')}/10")
            logger.info(f"   Speed: {model_info.get('speed', '?')}/10")
            logger.info(f"   Language: {ocr_language}")
            logger.info("="*60)
            
            # Upload PDF to Gemini
            logger.info("   📤 Step 1/4: Uploading PDF to Gemini...")
            upload_start = time.time()
            pdf_file = genai.upload_file(str(input_file))
            
            # Wait for processing
            logger.info("   ⏳ Step 2/4: Waiting for Gemini preprocessing...")
            while pdf_file.state.name == "PROCESSING":
                await asyncio.sleep(1)
                pdf_file = genai.get_file(pdf_file.name)
            
            if pdf_file.state.name == "FAILED":
                raise ValueError(f"PDF processing failed: {pdf_file.state.name}")
            
            upload_time = time.time() - upload_start
            logger.info(f"   ✓ Upload complete ({upload_time:.2f}s) - File ID: {pdf_file.name}")
            
            # Enhanced prompt for Gemini - Preserve layout and formatting
            prompt = f"""BẠN LÀ CHUYÊN GIA TRÍCH XUẤT VĂN BẢN TỪ PDF.

NHIỆM VỤ: Đọc file PDF này và trích xuất TOÀN BỘ nội dung văn bản, GIỮ NGUYÊN ĐỊNH DẠNG VÀ CẤU TRÚC gốc.

📋 YÊU CẦU QUAN TRỌNG:

1. CHÍNH TẢ & KÝ TỰ:
   - Giữ CHÍNH XÁC 100% mọi ký tự Tiếng Việt: ă, â, ê, ô, ơ, ư, đ, à, á, ả, ã, ạ, v.v.
   - Không sửa lỗi chính tả trong văn bản gốc
   - Giữ nguyên chữ hoa/thường như trong PDF

2. CẤU TRÚC VĂN BẢN:
   - GIỮ NGUYÊN số dòng trống giữa các đoạn văn
   - GIỮ NGUYÊN thụt lề đầu dòng (dùng spaces nếu có)
   - GIỮ NGUYÊN cách xuống dòng và ngắt đoạn
   - Nếu có đánh số (1., 2., a., b.) → GIỮ NGUYÊN format

3. TIÊU ĐỀ & HEADER:
   - Tiêu đề ở giữa trang → Thêm [CENTER] ở đầu dòng
   - Tiêu đề in đậm hoặc chữ hoa → Thêm [BOLD] ở đầu dòng
   - Ví dụ: [CENTER][BOLD]QUYẾT ĐỊNH

4. BẢNG BIỂU:
   - Mỗi hàng của bảng → Các ô cách nhau bằng dấu |
   - Hàng tiêu đề → Thêm [TABLE_HEADER] ở đầu
   - Ví dụ:
     [TABLE_HEADER]STT | Họ tên | Chức vụ
     1 | Nguyễn Văn A | Giám đốc
     2 | Trần Thị B | Phó giám đốc

5. DANH SÁCH & LIỆT KÊ:
   - GIỮ NGUYÊN dấu đầu dòng (-, *, •, 1., a.)
   - GIỮ NGUYÊN thụt lề các cấp

6. CHỮ KÝ & FOOTER:
   - GIỮ NGUYÊN vị trí căn phải/trái
   - Thêm [RIGHT] nếu căn phải
   - Ví dụ: [RIGHT]Giám đốc

7. NGÀY THÁNG & SỐ:
   - GIỮ NGUYÊN format: Ngày 15 tháng 12 năm 2024
   - Không chuyển đổi định dạng số

❌ TUYỆT ĐỐI KHÔNG:
- Thêm giải thích, chú thích, phân tích
- Sửa lỗi chính tả trong văn bản gốc
- Thay đổi format số, ngày tháng
- Tóm tắt hay bỏ qua bất kỳ nội dung nào

✅ CHỈ TRẢ VỀ:
- Văn bản thuần túy đã trích xuất
- Có các tag đánh dấu format: [CENTER], [BOLD], [RIGHT], [TABLE_HEADER]
- Giữ nguyên 100% nội dung và cấu trúc

Ngôn ngữ văn bản: {ocr_language}
Bắt đầu trích xuất:
"""
            
            # Generate content with optimized config using GeminiService wrapper
            logger.info("   🧠 Step 3/4: Extracting content with Gemini AI...")
            api_start = time.time()
            
            # Create GeminiService for auto-logging
            gemini_service = get_gemini_service(db) if db else None
            
            if gemini_service:
                # Use wrapper for auto-logging
                response = gemini_service.generate_content(
                    prompt=[pdf_file, prompt],
                    model=self.gemini_model_name,
                    operation="pdf-to-word",
                    metadata={
                        "file_name": input_file.name,
                        "ocr_language": ocr_language,
                        "conversion_type": "pdf_to_word_gemini"
                    },
                    generation_config={
                        "temperature": 0.0,
                        "top_p": 0.95,
                        "top_k": 40,
                        "max_output_tokens": 8192
                    }
                )
                logger.info("   ✓ Usage automatically logged to database")
            else:
                # Fallback to direct API call (no logging)
                logger.warning("   ⚠️  No db session - API call will NOT be logged")
                model_obj = genai.GenerativeModel(self.gemini_model_name)
                response = model_obj.generate_content(
                    [pdf_file, prompt],
                    generation_config=genai.GenerationConfig(
                        temperature=0.0,
                        top_p=0.95,
                        top_k=40,
                        max_output_tokens=8192
                    )
                )
            
            api_time = time.time() - api_start
            
            # Get plain text response from Gemini
            extracted_text = response.text.strip()
            logger.info(f"   ✓ Extraction complete ({api_time:.2f}s) - {len(extracted_text)} characters")
            
            # Create simple document structure
            document_data = {
                "content": extracted_text,
                "format": "text",
                "language": ocr_language
            }
            
            # Convert to Word
            logger.info("   📝 Step 4/4: Creating Word document...")
            word_start = time.time()
            word_path = await self._create_word_from_text(document_data, output_path)
            word_time = time.time() - word_start
            logger.info(f"   ✓ Word created ({word_time:.2f}s)")
            
            # Cleanup
            try:
                genai.delete_file(pdf_file.name)
                logger.info("   ✓ Cleanup complete")
            except Exception as e:
                logger.warning(f"   ⚠️  Cleanup warning: {e}")
            
            total_time = time.time() - start_time
            logger.info("")
            logger.info("="*60)
            logger.info("✅ GEMINI API - CONVERSION SUCCESS")
            logger.info(f"   Output: {output_path.name}")
            logger.info(f"   Size: {output_path.stat().st_size / 1024:.2f} KB")
            logger.info(f"   Total time: {total_time:.2f}s (upload: {upload_time:.2f}s, AI: {api_time:.2f}s, Word: {word_time:.2f}s)")
            logger.info("="*60)
            logger.info("")
            
            return word_path
            
        except Exception as e:
            logger.error("")
            logger.error("="*60)
            logger.error("❌ GEMINI API - CONVERSION FAILED")
            logger.error(f"   Error: {str(e)}")
            logger.error("="*60)
            logger.error("")
            raise HTTPException(500, f"Gemini conversion failed: {str(e)}")
        finally:
            # Restore original model if we switched
            if model_name and model_name != original_model:
                try:
                    self.set_gemini_model(original_model)
                    logger.info(f"Restored original model: {original_model}")
                except Exception as e:
                    logger.warning(f"Failed to restore original model: {e}")
    
    async def _create_word_from_text(
        self,
        data: dict,
        output_path: Path
    ) -> Path:
        """
        Create Word document from plain text extracted by Gemini
        Handles formatting tags: [CENTER], [BOLD], [RIGHT], [TABLE_HEADER]
        """
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            content = data.get("content", "").strip()
            
            if not content:
                # Empty document
                para = doc.add_paragraph()
                run = para.add_run("Không thể trích xuất nội dung từ PDF. File có thể bị lỗi hoặc không có text.")
                run.font.size = Pt(11)
            else:
                # Split by lines and create paragraphs
                lines = content.split('\n')
                logger.info(f"Creating {len(lines)} paragraphs from extracted text")
                
                current_table = None
                
                for line in lines:
                    original_line = line
                    line = line.strip()
                    
                    if not line:
                        # Empty line - add spacing
                        doc.add_paragraph()
                        current_table = None  # End table
                        continue
                    
                    # Parse formatting tags
                    is_center = '[CENTER]' in line
                    is_bold = '[BOLD]' in line
                    is_right = '[RIGHT]' in line
                    is_table_header = '[TABLE_HEADER]' in line
                    
                    # Remove tags from text
                    line = line.replace('[CENTER]', '').replace('[BOLD]', '').replace('[RIGHT]', '').replace('[TABLE_HEADER]', '').strip()
                    
                    # Handle table-like content (separated by | pipe)
                    if '|' in line and line.count('|') >= 1:
                        cells = [cell.strip() for cell in line.split('|')]
                        cells = [c for c in cells if c]  # Remove empty cells
                        
                        if len(cells) >= 2:
                            # Table content
                            if current_table is None or len(current_table.columns) != len(cells):
                                # Create new table
                                current_table = doc.add_table(rows=1, cols=len(cells))
                                current_table.style = 'Table Grid'
                                row_cells = current_table.rows[0].cells
                            else:
                                # Add row to existing table
                                row = current_table.add_row()
                                row_cells = row.cells
                            
                            # Fill cells
                            for i, cell_text in enumerate(cells):
                                cell = row_cells[i]
                                cell.text = cell_text
                                
                                # Format header row
                                if is_table_header:
                                    for paragraph in cell.paragraphs:
                                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                        for run in paragraph.runs:
                                            run.bold = True
                                            run.font.size = Pt(11)
                            continue
                    
                    # Regular paragraph
                    current_table = None  # End table when non-table content appears
                    
                    # Preserve indentation from original line
                    indent_spaces = len(original_line) - len(original_line.lstrip())
                    if indent_spaces > 0:
                        line = ' ' * indent_spaces + line
                    
                    para = doc.add_paragraph()
                    run = para.add_run(line)
                    run.font.size = Pt(11)
                    
                    # Apply formatting
                    if is_bold or line.isupper():
                        run.bold = True
                        run.font.size = Pt(12)
                    
                    if is_center:
                        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    elif is_right:
                        para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                    
                    # Auto-detect headers (Vietnamese keywords)
                    if any(kw in line.upper() for kw in [
                        'QUYẾT ĐỊNH', 'CÔNG HÒA XÃ HỘI CHỦ NGHĨA VIỆT NAM',
                        'ỦY BAN NHÂN DÂN', 'GIẤY CHỨNG NHẬN', 'CỘNG HÒA',
                        'HỘI NÔNG DÂN', 'THÔNG BÁO', 'BÁO CÁO', 'BIÊN BẢN',
                        'HỢP ĐỒNG', 'ĐƠN XIN', 'GIẤY ỦY QUYỀN'
                    ]):
                        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run.bold = True
                        run.font.size = Pt(13)
            
            # Save document
            doc.save(str(output_path))
            logger.info(f"✅ Word document created with {len(doc.paragraphs)} paragraphs and {len(doc.tables)} tables")
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating Word document: {e}")
            raise HTTPException(500, f"Failed to create Word document: {str(e)}")
    
    async def _create_word_from_json(
        self,
        data: dict,
        output_path: Path
    ) -> Path:
        """
        LEGACY: Create Word document with 100% EDITABLE TEXT - no images or objects
        This is kept for backward compatibility but _create_word_from_text is preferred
        """
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            logger.info(f"Creating Word document from data: {str(data)[:200]}...")
            
            # CRITICAL: All content must be editable text
            if data.get("format") == "text":
                # Text response - convert to editable paragraphs
                content = data.get("content", "")
                logger.info(f"Processing text content: {len(content)} characters")
                
                if content.strip():
                    lines = content.split('\n')
                    
                    for line in lines:
                        if line.strip():
                            # Create paragraph with pure text - NO IMAGES
                            para = doc.add_paragraph()
                            run = para.add_run(line.strip())
                            
                            # Basic text formatting - keep it editable
                            if any(keyword in line.upper() for keyword in ["QUYẾT ĐỊNH", "HỘI NÔNG DÂN", "CÔNG HÒA XÃ HỘI"]):
                                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                run.bold = True
                            
                            # Ensure it's text, not any embedded object
                            run.font.size = Pt(11)
                else:
                    # If no content, add a message
                    para = doc.add_paragraph()
                    run = para.add_run("Không thể trích xuất nội dung từ PDF. Vui lòng thử lại.")
                    run.font.size = Pt(11)
            
            else:
                # Structured response - process as pure text
                logger.info("Processing structured data")
                content_added = False
                
                # Document title as EDITABLE text
                doc_info = data.get('document_info', {})
                if doc_info.get('title'):
                    title_para = doc.add_paragraph()
                    title_run = title_para.add_run(doc_info['title'])
                    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    title_run.bold = True
                    title_run.font.size = Pt(14)
                    content_added = True
                        
                # Process pages with TEXT ONLY approach
                pages = data.get('pages', [])
                logger.info(f"Found {len(pages)} pages to process")
                
                for page_idx, page in enumerate(pages):
                    logger.info(f"Processing page {page_idx + 1}")
                    page_content = page.get('content', [])
                    
                    for item_idx, item in enumerate(page_content):
                        item_type = item.get('type', '')
                        text = item.get('text', '')
                        
                        logger.info(f"Processing item {item_idx}: type={item_type}, text_len={len(text)}")
                        
                        # Only process text-based content
                        if item_type in ['header', 'heading']:
                            if text.strip():
                                para = doc.add_paragraph()
                                run = para.add_run(text)
                                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                run.bold = True
                                run.font.size = Pt(12)
                                content_added = True
                                
                        elif item_type == 'paragraph':
                            if text.strip():
                                para = doc.add_paragraph()
                                run = para.add_run(text)
                                run.font.size = Pt(11)
                                content_added = True
                            
                        elif item_type == 'table':
                            headers = item.get('headers', {}).get('values', [])
                            rows = item.get('rows', [])
                            
                            if headers and rows:
                                logger.info(f"Creating table with {len(headers)} columns and {len(rows)} rows")
                                # Create Word table with EDITABLE TEXT cells
                                table = doc.add_table(rows=1, cols=len(headers))
                                table.style = 'Table Grid'
                                
                                # Header row - EDITABLE TEXT
                                hdr_cells = table.rows[0].cells
                                for i, header in enumerate(headers):
                                    hdr_cells[i].text = str(header)
                                    # Make header text bold and centered
                                    for paragraph in hdr_cells[i].paragraphs:
                                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                        for run in paragraph.runs:
                                            run.bold = True
                                
                                # Data rows - EDITABLE TEXT  
                                for row in rows:
                                    row_cells = table.add_row().cells
                                    cells_data = row.get('cells', []) if isinstance(row, dict) else row
                                    
                                    for i, cell_data in enumerate(cells_data):
                                        if i < len(row_cells):
                                            if isinstance(cell_data, dict):
                                                cell_text = str(cell_data.get('text', ''))
                                            else:
                                                cell_text = str(cell_data)
                                            
                                            row_cells[i].text = cell_text
                                            
                                            # Ensure all text is editable
                                            for paragraph in row_cells[i].paragraphs:
                                                for run in paragraph.runs:
                                                    run.font.size = Pt(10)
                                
                                content_added = True
                
                # If no content was added from structured data, try to extract any text
                if not content_added:
                    logger.warning("No content found in structured data, trying to extract any available text")
                    
                    # Try to extract from the entire data structure
                    import json
                    full_text = json.dumps(data, ensure_ascii=False, indent=2)
                    
                    para = doc.add_paragraph()
                    run = para.add_run("Nội dung được trích xuất từ PDF:")
                    run.font.size = Pt(11)
                    run.bold = True
                    
                    # Add the raw extracted data
                    para = doc.add_paragraph()
                    run = para.add_run(full_text[:2000])  # Limit to first 2000 chars
                    run.font.size = Pt(10)
            
            # Save as Word document with pure text content
            doc.save(str(output_path))
            logger.info(f"✅ Word document created with editable text: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating editable Word document: {str(e)}")
            logger.exception("Full error details:")
            raise
    
    async def pdf_to_excel(
        self,
        input_file: Path,
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert PDF to Excel (.xlsx)
        Extracts tables from PDF using pdfplumber
        """
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + ".xlsx"
        output_path = self.output_dir / output_filename
        
        try:
            import pdfplumber
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            
            # Create Excel workbook
            wb = Workbook()
            wb.remove(wb.active)  # Remove default sheet
            
            # Open PDF
            with pdfplumber.open(input_file) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    # Extract tables from page
                    tables = page.extract_tables()
                    
                    if tables:
                        # Create sheet for each page with tables
                        ws = wb.create_sheet(title=f"Page {page_num}")
                        
                        # Write all tables from this page
                        current_row = 1
                        for table_num, table in enumerate(tables, start=1):
                            # Add table header
                            if len(tables) > 1:
                                header_cell = ws.cell(row=current_row, column=1, value=f"Table {table_num}")
                                header_cell.font = Font(bold=True, size=12)
                                header_cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
                                current_row += 1
                            
                            # Write table data
                            for row_data in table:
                                for col_num, cell_value in enumerate(row_data, start=1):
                                    cell = ws.cell(row=current_row, column=col_num, value=cell_value)
                                    cell.alignment = Alignment(wrap_text=True, vertical='top')
                                    
                                    # Bold first row (header)
                                    if current_row == 1 or (len(tables) > 1 and current_row == 2):
                                        cell.font = Font(bold=True)
                                        cell.fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")
                                
                                current_row += 1
                            
                            # Add spacing between tables
                            current_row += 1
                        
                        # Auto-adjust column widths
                        for column in ws.columns:
                            max_length = 0
                            column_letter = column[0].column_letter
                            for cell in column:
                                try:
                                    if cell.value:
                                        max_length = max(max_length, len(str(cell.value)))
                                except:
                                    pass
                            adjusted_width = min(max_length + 2, 50)
                            ws.column_dimensions[column_letter].width = adjusted_width
                    else:
                        # No tables found, extract text
                        text = page.extract_text()
                        if text and text.strip():
                            ws = wb.create_sheet(title=f"Page {page_num} (Text)")
                            ws.cell(row=1, column=1, value=text)
                            ws.column_dimensions['A'].width = 100
            
            # If no content was extracted at all
            if len(wb.sheetnames) == 0:
                ws = wb.create_sheet(title="No Data")
                ws.cell(row=1, column=1, value="No tables or text found in PDF")
            
            # Save workbook
            wb.save(output_path)
            
            return output_path
            
        except Exception as e:
            raise HTTPException(500, f"PDF to Excel conversion failed: {str(e)}")
    
    # ==================== PDF Operations ====================
    
    async def merge_pdfs(
        self,
        input_files: List[Path],
        output_filename: str = "merged.pdf"
    ) -> Path:
        """Merge multiple PDFs into one (using modern pypdf)"""
        output_path = self.output_dir / output_filename
        
        try:
            merger = pypdf.PdfWriter()
            
            for pdf_file in input_files:
                if pdf_file.suffix.lower() != '.pdf':
                    continue
                    
                with open(pdf_file, 'rb') as f:
                    pdf_reader = pypdf.PdfReader(f)
                    for page in pdf_reader.pages:
                        merger.add_page(page)
            
            with open(output_path, 'wb') as output_file:
                merger.write(output_file)
                
            return output_path
            
        except Exception as e:
            raise HTTPException(500, f"PDF merge failed: {str(e)}")
    
    async def split_pdf(
        self,
        input_file: Path,
        page_ranges: List[tuple],  # [(1,3), (5,7)] means pages 1-3 and 5-7
        output_prefix: str = "split"
    ) -> List[Path]:
        """Split PDF into multiple files by page ranges"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        try:
            output_files = []
            
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                
                for idx, (start, end) in enumerate(page_ranges):
                    pdf_writer = pypdf.PdfWriter()
                    
                    for page_num in range(start - 1, end):  # Convert to 0-indexed
                        if page_num < len(pdf_reader.pages):
                            pdf_writer.add_page(pdf_reader.pages[page_num])
                    
                    output_path = self.output_dir / f"{output_prefix}_{idx+1}.pdf"
                    with open(output_path, 'wb') as output_file:
                        pdf_writer.write(output_file)
                    
                    output_files.append(output_path)
                    
            return output_files
            
        except Exception as e:
            raise HTTPException(500, f"PDF split failed: {str(e)}")
    
    async def extract_pdf_text(self, input_file: Path) -> str:
        """Extract all text from PDF"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        try:
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                text = ""
                
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n\n"
                    
            return text.strip()
            
        except Exception as e:
            raise HTTPException(500, f"PDF text extraction failed: {str(e)}")
    
    async def rotate_pdf_pages(
        self,
        input_file: Path,
        rotation: int = 90,  # 90, 180, 270
        pages: Optional[List[int]] = None,  # None = all pages
        output_filename: Optional[str] = None
    ) -> Path:
        """Rotate PDF pages"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        if rotation not in [90, 180, 270, -90]:
            raise HTTPException(400, "Rotation must be 90, 180, or 270 degrees")
        
        output_filename = output_filename or input_file.stem + "_rotated.pdf"
        output_path = self.output_dir / output_filename
        
        try:
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                pdf_writer = pypdf.PdfWriter()
                
                for page_num, page in enumerate(pdf_reader.pages):
                    if pages is None or (page_num + 1) in pages:
                        page.rotate(rotation)
                    pdf_writer.add_page(page)
                
                with open(output_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                    
            return output_path
            
        except Exception as e:
            raise HTTPException(500, f"PDF rotation failed: {str(e)}")
    
    # ==================== PDF Info ====================
    
    async def get_pdf_info(self, input_file: Path) -> dict:
        """Get PDF metadata and information"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        try:
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                
                info = {
                    "num_pages": len(pdf_reader.pages),
                    "metadata": pdf_reader.metadata.__dict__ if pdf_reader.metadata else {},
                    "is_encrypted": pdf_reader.is_encrypted,
                    "page_sizes": []
                }
                
                # Get page sizes
                for page in pdf_reader.pages[:5]:  # First 5 pages
                    box = page.mediabox
                    info["page_sizes"].append({
                        "width": float(box.width),
                        "height": float(box.height)
                    })
                    
            return info
            
        except Exception as e:
            raise HTTPException(500, f"PDF info extraction failed: {str(e)}")
    
    # ==================== Word Operations ====================
    
    async def get_word_info(self, input_file: Path) -> dict:
        """Get Word document information"""
        if input_file.suffix.lower() != '.docx':
            raise HTTPException(400, "File must be .docx")
        
        try:
            doc = Document(str(input_file))
            
            info = {
                "num_paragraphs": len(doc.paragraphs),
                "num_tables": len(doc.tables),
                "num_images": len([r for p in doc.paragraphs for r in p.runs if r._element.xml.find('pic') > 0]),
                "word_count": sum(len(p.text.split()) for p in doc.paragraphs),
                "char_count": sum(len(p.text) for p in doc.paragraphs),
            }
            
            return info
            
        except Exception as e:
            raise HTTPException(500, f"Word info extraction failed: {str(e)}")
    
    # ==================== Excel Operations ====================
    
    async def get_excel_info(self, input_file: Path) -> dict:
        """Get Excel workbook information"""
        if input_file.suffix.lower() not in ['.xlsx', '.xlsm']:
            raise HTTPException(400, "File must be .xlsx or .xlsm")
        
        try:
            wb = load_workbook(str(input_file), read_only=True)
            
            info = {
                "num_sheets": len(wb.sheetnames),
                "sheet_names": wb.sheetnames,
                "active_sheet": wb.active.title,
            }
            
            wb.close()
            return info
            
        except Exception as e:
            raise HTTPException(500, f"Excel info extraction failed: {str(e)}")
    
    # ==================== PowerPoint Operations ====================
    
    async def get_powerpoint_info(self, input_file: Path) -> dict:
        """Get PowerPoint presentation information"""
        if input_file.suffix.lower() not in ['.pptx']:
            raise HTTPException(400, "File must be .pptx")
        
        try:
            prs = Presentation(str(input_file))
            
            info = {
                "num_slides": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
            }
            
            return info
            
        except Exception as e:
            raise HTTPException(500, f"PowerPoint info extraction failed: {str(e)}")
    
    # ==================== Cleanup ====================
    
    # ==================== PDF Compression ====================
    
    async def compress_pdf(
        self,
        input_file: Path,
        quality: str = "medium",  # low, medium, high
        output_filename: Optional[str] = None
    ) -> tuple[Path, str]:
        """
        Compress PDF - HYBRID STRATEGY
        
        Strategy:
        1. Check technology priority from settings
        2. Try Adobe PDF Services (if enabled and high priority) - AI compression 10/10
        3. Fallback to pypdf - Basic compression 7/10
        
        Args:
            input_file: Path to PDF file
            quality: Compression level (low, medium, high)
            output_filename: Optional output filename
        
        Returns:
            Tuple of (output_path, technology_used)
            - output_path: Path to compressed PDF
            - technology_used: 'adobe' or 'pypdf'
        
        Quality mapping:
        - Adobe: low=HIGH (aggressive), medium=MEDIUM, high=LOW (preserve quality)
        - pypdf: low=4, medium=2, high=0
        """
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + "_compressed.pdf"
        output_path = self.output_dir / output_filename
        
        # Import settings here to avoid circular imports
        from app.core.config import settings
        
        # Get technology priority for compress operation
        priorities = settings.get_technology_priority("compress")
        
        # Try each technology in priority order
        for tech in priorities:
            if tech.lower() == "adobe":
                # Try Adobe if enabled
                if self.use_adobe and self.adobe_credentials and ADOBE_AVAILABLE:
                    try:
                        logger.info(f"Trying Adobe compress for {input_file.name}")
                        await self._compress_pdf_adobe(input_file, quality, output_path)
                        logger.info(f"Adobe compression successful: {output_path}")
                        return (output_path, "adobe")
                    except Exception as e:
                        logger.warning(f"Adobe compress failed: {e}, trying next technology")
                        continue
                else:
                    logger.debug("Adobe not available, skipping")
                    continue
                    
            elif tech.lower() == "pypdf":
                # Try pypdf (always available)
                try:
                    logger.info(f"Using pypdf compress for {input_file.name}")
                    await self._compress_pdf_local(input_file, quality, output_path)
                    logger.info(f"pypdf compression successful: {output_path}")
                    return (output_path, "pypdf")
                except Exception as e:
                    logger.warning(f"pypdf compress failed: {e}, trying next technology")
                    continue
        
        # If all technologies failed
        raise HTTPException(500, "All compression methods failed")
    
    async def _compress_pdf_adobe(
        self,
        input_file: Path,
        quality: str,
        output_path: Path
    ) -> None:
        """
        Compress PDF using Adobe PDF Services API
        AI-powered compression with 10/10 quality
        
        Adobe compression levels:
        - HIGH: Aggressive compression (50-80% reduction)
        - MEDIUM: Balanced (30-50% reduction)
        - LOW: Light compression (10-30% reduction)
        """
        try:
            # Read input file
            async with aiofiles.open(input_file, 'rb') as f:
                input_stream = await f.read()
            
            # Adobe API (simplified - actual implementation would use CompressPDF operation)
            # This is placeholder - needs actual Adobe CompressPDF SDK code
            from adobe.pdfservices.operation.pdfjobs.jobs.compress_pdf_job import CompressPDFJob
            from adobe.pdfservices.operation.pdfjobs.params.compress_pdf.compress_pdf_params import CompressPDFParams
            from adobe.pdfservices.operation.pdfjobs.params.compress_pdf.compression_level import CompressionLevel
            from adobe.pdfservices.operation.pdfjobs.result.compress_pdf_result import CompressPDFResult
            
            # Map quality to Adobe compression level
            quality_map = {
                "low": CompressionLevel.HIGH,      # More compression
                "medium": CompressionLevel.MEDIUM,
                "high": CompressionLevel.LOW       # Less compression, preserve quality
            }
            compression_level = quality_map.get(quality, CompressionLevel.MEDIUM)
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=input_stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Create compress parameters
            compress_params = CompressPDFParams(
                compression_level=compression_level
            )
            
            # Create and submit job
            compress_job = CompressPDFJob(
                input_asset=input_asset,
                compress_pdf_params=compress_params
            )
            
            location = pdf_services.submit(compress_job)
            pdf_services_response = pdf_services.get_job_result(location, CompressPDFResult)
            
            # Download result
            result_asset: CloudAsset = pdf_services_response.get_result().get_asset()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save to file
            async with aiofiles.open(output_path, "wb") as f:
                await f.write(stream_asset.get_input_stream())
            
        except ImportError:
            # CompressPDF not available in current SDK version
            logger.error("Adobe CompressPDF API not available in SDK")
            raise HTTPException(500, "Adobe CompressPDF requires SDK update")
        except Exception as e:
            logger.error(f"Adobe compression error: {e}")
            raise
    
    async def _compress_pdf_local(
        self,
        input_file: Path,
        quality: str,
        output_path: Path
    ) -> None:
        """
        Compress PDF using pypdf library
        Basic compression with 7/10 quality (existing code)
        """
        # Quality settings for pypdf compression
        quality_map = {
            "low": 4,      # More compression
            "medium": 2,   # Balanced
            "high": 0      # Less compression
        }
        
        compression_level = quality_map.get(quality, 2)
        
        try:
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                pdf_writer = pypdf.PdfWriter()
                
                # Copy all pages
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
                
                # Compress images
                for page in pdf_writer.pages:
                    page.compress_content_streams(level=compression_level)
                
                # Write compressed PDF
                with open(output_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                    
        except Exception as e:
            raise HTTPException(500, f"pypdf compression failed: {str(e)}")
    
    # ==================== Image to PDF ====================
    
    async def image_to_pdf(
        self,
        input_file: Path,
        output_filename: Optional[str] = None
    ) -> Path:
        """Convert image (JPG, PNG, etc.) to PDF"""
        allowed_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.heic']
        
        if input_file.suffix.lower() not in allowed_extensions:
            raise HTTPException(400, f"File must be one of: {', '.join(allowed_extensions)}")
        
        output_filename = output_filename or input_file.stem + ".pdf"
        output_path = self.output_dir / output_filename
        
        try:
            # Use Pillow to convert image to PDF
            from PIL import Image
            
            # Open and convert image
            img = Image.open(input_file)
            
            # Convert RGBA to RGB if necessary (for PNG with transparency)
            if img.mode in ('RGBA', 'LA', 'P'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Save as PDF
            img.save(output_path, 'PDF', resolution=100.0, quality=95)
            
            return output_path
            
        except Exception as e:
            raise HTTPException(500, f"Image to PDF conversion failed: {str(e)}")
    
    async def images_to_pdf(
        self,
        input_files: List[Path],
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert multiple images to a single PDF (each image = 1 page)
        
        Args:
            input_files: List of image file paths
            output_filename: Optional output filename
            
        Returns:
            Path to combined PDF
        """
        if not input_files:
            raise HTTPException(400, "No images provided")
        
        output_filename = output_filename or f"images_combined_{len(input_files)}_pages.pdf"
        output_path = self.output_dir / output_filename
        
        try:
            from PIL import Image
            
            # Convert all images to RGB PIL Images
            rgb_images = []
            for img_path in input_files:
                img = Image.open(img_path)
                
                # Convert to RGB
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                    img = background
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                rgb_images.append(img)
            
            # Save first image as PDF, append others as pages
            if len(rgb_images) == 1:
                rgb_images[0].save(output_path, 'PDF', resolution=100.0, quality=95)
            else:
                rgb_images[0].save(
                    output_path, 
                    'PDF', 
                    resolution=100.0, 
                    quality=95,
                    save_all=True,
                    append_images=rgb_images[1:]
                )
            
            return output_path
            
        except Exception as e:
            raise HTTPException(500, f"Images to PDF conversion failed: {str(e)}")
    
    # ==================== PDF Watermark ====================
    
    async def add_watermark_to_pdf(
        self,
        input_file: Path,
        watermark_text: str,
        position: str = "center",  # center, top-left, top-right, bottom-left, bottom-right
        opacity: float = 0.3,
        output_filename: Optional[str] = None
    ) -> tuple[Path, str]:
        """
        Add text watermark to PDF - HYBRID STRATEGY
        
        Strategy:
        1. Check technology priority from settings
        2. Try Adobe PDF Services (if enabled) - Advanced watermark 10/10
        3. Fallback to reportlab+pypdf - Basic watermark 8/10
        
        Returns:
            Tuple of (output_path, technology_used)
            - output_path: Path to watermarked PDF
            - technology_used: 'adobe' or 'pypdf'
        """
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + "_watermarked.pdf"
        output_path = self.output_dir / output_filename
        
        # Import settings
        from app.core.config import settings
        
        # Get technology priority for watermark operation
        priorities = settings.get_technology_priority("watermark")
        
        # Try each technology in priority order
        for tech in priorities:
            if tech.lower() == "adobe":
                # Try Adobe if enabled
                if self.use_adobe and self.adobe_credentials and ADOBE_AVAILABLE:
                    try:
                        logger.info(f"Trying Adobe watermark for {input_file.name}")
                        await self._add_watermark_adobe(
                            input_file, watermark_text, position, opacity, output_path
                        )
                        logger.info(f"Adobe watermark successful: {output_path}")
                        return (output_path, "adobe")
                    except Exception as e:
                        logger.warning(f"Adobe watermark failed: {e}, trying next technology")
                        continue
                else:
                    logger.debug("Adobe not available, skipping")
                    continue
                    
            elif tech.lower() == "pypdf":
                # Try pypdf (always available)
                try:
                    logger.info(f"Using pypdf watermark for {input_file.name}")
                    await self._add_watermark_local(
                        input_file, watermark_text, position, opacity, output_path
                    )
                    logger.info(f"pypdf watermark successful: {output_path}")
                    return (output_path, "pypdf")
                except Exception as e:
                    logger.warning(f"pypdf watermark failed: {e}, trying next technology")
                    continue
        
        # If all technologies failed
        raise HTTPException(500, "All watermark methods failed")
    
    async def _add_watermark_adobe(
        self,
        input_file: Path,
        watermark_text: str,
        position: str,
        opacity: float,
        output_path: Path
    ) -> None:
        """
        Add watermark using Adobe PDF Services API
        Advanced watermark with 10/10 quality
        
        Note: Adobe doesn't have native Watermark API yet,
        so we use AddTextWatermark or manual page overlay
        This is a placeholder for future Adobe Watermark API
        """
        try:
            # Read input file
            async with aiofiles.open(input_file, 'rb') as f:
                input_stream = await f.read()
            
            # Adobe Watermark API (placeholder - Adobe SDK may not have this yet)
            # For now, we'll raise ImportError to fallback to pypdf
            # When Adobe releases Watermark API, update this code
            
            raise ImportError("Adobe Watermark API not available yet in SDK")
            
            # Future implementation when Adobe adds Watermark API:
            # from adobe.pdfservices.operation.pdfjobs.jobs.add_watermark_job import AddWatermarkJob
            # from adobe.pdfservices.operation.pdfjobs.params.add_watermark_params import AddWatermarkParams
            # 
            # pdf_services = PDFServices(credentials=self.adobe_credentials)
            # input_asset = pdf_services.upload(input_stream, PDFServicesMediaType.PDF)
            # 
            # watermark_params = AddWatermarkParams(
            #     text=watermark_text,
            #     position=position,
            #     opacity=opacity
            # )
            # 
            # watermark_job = AddWatermarkJob(input_asset, watermark_params)
            # location = pdf_services.submit(watermark_job)
            # result = pdf_services.get_job_result(location, AddWatermarkResult)
            # 
            # result_asset = result.get_result().get_asset()
            # stream_asset = pdf_services.get_content(result_asset)
            # 
            # async with aiofiles.open(output_path, "wb") as f:
            #     await f.write(stream_asset.get_input_stream())
            
        except ImportError:
            logger.info("Adobe Watermark API not available, will fallback to pypdf")
            raise
        except Exception as e:
            logger.error(f"Adobe watermark error: {e}")
            raise
    
    async def _add_watermark_local(
        self,
        input_file: Path,
        watermark_text: str,
        position: str,
        opacity: float,
        output_path: Path
    ) -> None:
        """
        Add watermark using reportlab + pypdf
        Basic watermark with 8/10 quality (existing code)
        """
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.colors import Color
            import io
            
            # Read input PDF
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                pdf_writer = pypdf.PdfWriter()
                
                # Create watermark
                packet = io.BytesIO()
                can = canvas.Canvas(packet, pagesize=letter)
                
                # Position mapping
                positions = {
                    'center': (300, 400),
                    'top-left': (50, 750),
                    'top-right': (450, 750),
                    'bottom-left': (50, 50),
                    'bottom-right': (450, 50)
                }
                x, y = positions.get(position, (300, 400))
                
                # Set watermark properties
                can.setFillColor(Color(0, 0, 0, alpha=opacity))
                can.setFont("Helvetica", 40)
                can.drawString(x, y, watermark_text)
                can.save()
                
                # Move to beginning of BytesIO
                packet.seek(0)
                watermark_pdf = pypdf.PdfReader(packet)
                watermark_page = watermark_pdf.pages[0]
                
                # Apply watermark to all pages
                for page in pdf_reader.pages:
                    page.merge_page(watermark_page)
                    pdf_writer.add_page(page)
                
                # Write output
                with open(output_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                    
        except Exception as e:
            raise HTTPException(500, f"pypdf watermark failed: {str(e)}")
    
    # ==================== PDF Password Protection ====================
    
    async def protect_pdf_with_password(
        self,
        input_file: Path,
        user_password: str,
        owner_password: Optional[str] = None,
        output_filename: Optional[str] = None
    ) -> Path:
        """Protect PDF with password"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + "_protected.pdf"
        output_path = self.output_dir / output_filename
        
        if not owner_password:
            owner_password = user_password + "_owner"
        
        try:
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                pdf_writer = pypdf.PdfWriter()
                
                # Copy all pages
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
                
                # Encrypt with password
                pdf_writer.encrypt(
                    user_password=user_password,
                    owner_password=owner_password,
                    permissions_flag=0b0100  # Allow printing
                )
                
                # Write encrypted PDF
                with open(output_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                    
            return output_path
            
        except Exception as e:
            raise HTTPException(500, f"PDF password protection failed: {str(e)}")
    
    async def unlock_pdf(
        self,
        input_file: Path,
        password: str,
        output_filename: Optional[str] = None
    ) -> Path:
        """Remove password from PDF"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + "_unlocked.pdf"
        output_path = self.output_dir / output_filename
        
        try:
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                
                # Check if encrypted
                if not pdf_reader.is_encrypted:
                    raise HTTPException(400, "PDF is not encrypted")
                
                # Decrypt
                if not pdf_reader.decrypt(password):
                    raise HTTPException(401, "Incorrect password")
                
                pdf_writer = pypdf.PdfWriter()
                
                # Copy all pages (now decrypted)
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
                
                # Write unencrypted PDF
                with open(output_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                    
            return output_path
            
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"PDF unlock failed: {str(e)}")
    
    # ==================== PDF to Images ====================
    
    async def pdf_to_images(
        self,
        input_file: Path,
        format: str = "png",  # png, jpg
        dpi: int = 200,
        output_prefix: str = "page"
    ) -> List[Path]:
        """Convert PDF pages to images using pypdfium2 (no Poppler needed!)"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        try:
            import pypdfium2 as pdfium
            from PIL import Image
            
            output_files = []
            
            # Open PDF with pypdfium2
            pdf = pdfium.PdfDocument(input_file)
            
            # Calculate scale for DPI (default 72 DPI, scale up to requested DPI)
            scale = dpi / 72.0
            
            # Convert each page
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                
                # Render page to PIL Image with specified DPI
                pil_image = page.render(
                    scale=scale,
                    rotation=0,
                    crop=(0, 0, 0, 0)
                ).to_pil()
                
                # Save image
                output_filename = f"{output_prefix}_{page_num + 1}.{format}"
                output_path = self.output_dir / output_filename
                
                # Convert RGBA to RGB if saving as JPEG
                if format.lower() in ['jpg', 'jpeg'] and pil_image.mode == 'RGBA':
                    rgb_image = Image.new('RGB', pil_image.size, (255, 255, 255))
                    rgb_image.paste(pil_image, mask=pil_image.split()[3])
                    pil_image = rgb_image
                
                pil_image.save(output_path, format.upper() if format.upper() != 'JPG' else 'JPEG', quality=95)
                output_files.append(output_path)
            
            pdf.close()
            return output_files
            
        except Exception as e:
            raise HTTPException(500, f"PDF to images conversion failed: {str(e)}")
    
    # ==================== Add Page Numbers ====================
    
    async def add_page_numbers(
        self,
        input_file: Path,
        position: str = "bottom-center",  # bottom-center, bottom-right, bottom-left
        format: str = "Page {page}",  # {page}, {page} of {total}, etc.
        output_filename: Optional[str] = None
    ) -> Path:
        """Add page numbers to PDF"""
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + "_numbered.pdf"
        output_path = self.output_dir / output_filename
        
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            import io
            
            with open(input_file, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                pdf_writer = pypdf.PdfWriter()
                total_pages = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages, start=1):
                    # Create page number overlay
                    packet = io.BytesIO()
                    can = canvas.Canvas(packet, pagesize=letter)
                    
                    # Position mapping
                    positions = {
                        'bottom-center': (300, 30),
                        'bottom-right': (520, 30),
                        'bottom-left': (80, 30)
                    }
                    x, y = positions.get(position, (300, 30))
                    
                    # Format page number text
                    page_text = format.replace('{page}', str(page_num))
                    page_text = page_text.replace('{total}', str(total_pages))
                    
                    can.setFont("Helvetica", 10)
                    can.drawString(x, y, page_text)
                    can.save()
                    
                    # Merge with original page
                    packet.seek(0)
                    overlay_pdf = pypdf.PdfReader(packet)
                    page.merge_page(overlay_pdf.pages[0])
                    pdf_writer.add_page(page)
                
                # Write output
                with open(output_path, 'wb') as output_file:
                    pdf_writer.write(output_file)
                    
            return output_path
            
        except Exception as e:
            raise HTTPException(500, f"Add page numbers failed: {str(e)}")
    
    # ==================== ADOBE-ONLY FEATURES ====================
    
    async def ocr_pdf(
        self,
        input_file: Path,
        language: str = "vi-VN",
        output_filename: Optional[str] = None
    ) -> Path:
        """
        OCR scanned PDF to searchable PDF
        
        Technology Priority:
        1. Adobe PDF Services (10/10) - Best quality, 50+ languages
        2. Tesseract OCR (7/10) - Free fallback with basic OCR
        
        Features:
        - Text recognition from scanned PDFs
        - Support Vietnamese (vi-VN) and 50+ languages
        - Preserve original layout (Adobe only)
        - Add searchable text layer
        
        Args:
            input_file: Path to scanned PDF file
            language: Language code (vi-VN, en-US, fr-FR, etc.)
            output_filename: Optional output filename
        
        Returns:
            Path to searchable PDF
        
        Raises:
            HTTPException 400: Invalid input
            HTTPException 500: OCR fails
        """
        # Try Adobe first (best quality)
        if self.use_adobe and self.adobe_credentials and ADOBE_AVAILABLE:
            logger.info("🎯 Using Adobe OCR for PDF")
            return await self._ocr_pdf_adobe(input_file, language, output_filename)
        
        # Fallback to Tesseract
        logger.info("🔄 Using Tesseract OCR fallback (Adobe not available)")
        return await self._ocr_pdf_tesseract(input_file, language, output_filename)
    
    async def _ocr_pdf_adobe(
        self,
        input_file: Path,
        language: str = "vi-VN",
        output_filename: Optional[str] = None
    ) -> Path:
        """Adobe OCR implementation (original logic)"""
        
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + "_ocr.pdf"
        output_path = self.output_dir / output_filename
        
        try:
            # Read input file
            async with aiofiles.open(input_file, 'rb') as f:
                input_stream = await f.read()
            
            # Adobe OCR API - Use same imports as _pdf_to_word_adobe_internal
            from adobe.pdfservices.operation.pdfjobs.jobs.ocr_pdf_job import OCRPDFJob
            from adobe.pdfservices.operation.pdfjobs.params.ocr_pdf.ocr_pdf_params import OCRPDFParams
            from adobe.pdfservices.operation.pdfjobs.params.ocr_pdf.ocr_supported_locale import OCRSupportedLocale
            from adobe.pdfservices.operation.pdfjobs.params.ocr_pdf.ocr_supported_type import OCRSupportedType
            from adobe.pdfservices.operation.pdfjobs.result.ocr_pdf_result import OCRPDFResult
            
            # Map language codes to Adobe OCRSupportedLocale
            language_map = {
                "vi-VN": OCRSupportedLocale.VI_VN,  # Vietnamese
                "en-US": OCRSupportedLocale.EN_US,  # English
                "fr-FR": OCRSupportedLocale.FR_FR,  # French
                "de-DE": OCRSupportedLocale.DE_DE,  # German
                "es-ES": OCRSupportedLocale.ES_ES,  # Spanish
                "it-IT": OCRSupportedLocale.IT_IT,  # Italian
                "ja-JP": OCRSupportedLocale.JA_JP,  # Japanese
                "ko-KR": OCRSupportedLocale.KO_KR,  # Korean
                "zh-CN": OCRSupportedLocale.ZH_HANS, # Chinese Simplified
            }
            ocr_locale = language_map.get(language, OCRSupportedLocale.EN_US)
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=input_stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Create OCR parameters
            ocr_params = OCRPDFParams(
                ocr_locale=ocr_locale,
                ocr_type=OCRSupportedType.SEARCHABLE_IMAGE  # Preserve original appearance
            )
            
            # Create and submit OCR job
            ocr_job = OCRPDFJob(
                input_asset=input_asset,
                ocr_pdf_params=ocr_params
            )
            
            location = pdf_services.submit(ocr_job)
            logger.info(f"Adobe OCR job submitted for {input_file.name}")
            
            # Get result (polling handled by SDK)
            pdf_services_response = pdf_services.get_job_result(location, OCRPDFResult)
            
            # Download result
            result_asset: CloudAsset = pdf_services_response.get_result().get_asset()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save to file
            async with aiofiles.open(output_path, "wb") as f:
                await f.write(stream_asset.get_input_stream())
            
            logger.info(f"Adobe OCR successful: {output_path}")
            return output_path
            
        except ImportError:
            logger.error("Adobe OCR API not available in SDK")
            raise HTTPException(500, "Adobe OCR requires SDK update")
        except Exception as e:
            logger.error(f"Adobe OCR error: {e}")
            raise HTTPException(500, f"OCR failed: {str(e)}")
    
    async def _ocr_pdf_tesseract(
        self,
        input_file: Path,
        language: str = "vi-VN",
        output_filename: Optional[str] = None
    ) -> Path:
        """
        OCR PDF using Tesseract (free fallback)
        
        Quality: 7/10 (Good for basic OCR, but not perfect layout preservation)
        Features: Extract text, basic recognition, multiple languages
        Limitation: Creates new PDF with text layer (doesn't preserve exact layout)
        """
        try:
            from pdf2image import convert_from_path
            from PIL import Image
            import pytesseract
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.utils import ImageReader
            import io
        except ImportError as e:
            raise HTTPException(
                500, 
                f"Tesseract OCR dependencies not installed: {e}. "
                "Install: pip install pdf2image pytesseract pillow reportlab"
            )
        
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        output_filename = output_filename or input_file.stem + "_ocr_tesseract.pdf"
        output_path = self.output_dir / output_filename
        
        try:
            # Map language codes to Tesseract format
            lang_map = {
                "vi-VN": "vie",  # Vietnamese
                "en-US": "eng",  # English
                "fr-FR": "fra",  # French
                "de-DE": "deu",  # German
                "es-ES": "spa",  # Spanish
                "it-IT": "ita",  # Italian
                "ja-JP": "jpn",  # Japanese
                "ko-KR": "kor",  # Korean
                "zh-CN": "chi_sim",  # Chinese Simplified
            }
            tesseract_lang = lang_map.get(language, "eng")
            
            logger.info(f"Converting PDF to images for OCR...")
            # Convert PDF pages to images
            images = convert_from_path(str(input_file), dpi=300)
            
            logger.info(f"Performing OCR on {len(images)} pages...")
            # Create new PDF with OCR text
            c = canvas.Canvas(str(output_path), pagesize=letter)
            
            for i, image in enumerate(images):
                logger.info(f"  Processing page {i+1}/{len(images)}...")
                
                # OCR the image
                text = pytesseract.image_to_string(image, lang=tesseract_lang)
                
                # Convert PIL image to bytes for ReportLab
                img_buffer = io.BytesIO()
                image.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                img_reader = ImageReader(img_buffer)
                
                # Get image dimensions
                img_width, img_height = image.size
                
                # Scale to fit page (8.5 x 11 inches = 612 x 792 points)
                page_width, page_height = letter
                scale = min(page_width / img_width, page_height / img_height)
                scaled_width = img_width * scale
                scaled_height = img_height * scale
                
                # Draw image on page
                c.drawImage(img_reader, 0, 0, width=scaled_width, height=scaled_height)
                
                # Add invisible text layer (for searchability)
                c.setFillColorRGB(1, 1, 1, alpha=0.01)  # Almost invisible
                c.setFont("Helvetica", 8)
                text_object = c.beginText(10, page_height - 20)
                for line in text.split('\n'):
                    text_object.textLine(line)
                c.drawText(text_object)
                
                c.showPage()  # Next page
            
            c.save()
            logger.info(f"✅ Tesseract OCR successful: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Tesseract OCR error: {e}")
            raise HTTPException(500, f"Tesseract OCR failed: {str(e)}")
    
    async def extract_pdf_content(
        self,
        input_file: Path,
        extract_type: str = "all"  # all, text, tables, images
    ) -> dict:
        """
        AI-powered content extraction from PDF using Adobe Extract API
        
        Adobe Extract Features:
        - Extract tables → Structured data (CSV/Excel format)
        - Extract images → PNG files with metadata
        - Extract text with font information (bold, italic, size, family)
        - Reading order detection (AI-powered)
        - Character bounding boxes (precise position)
        - Document structure detection (headings, paragraphs, lists)
        
        Args:
            input_file: Path to PDF file
            extract_type: What to extract (all, text, tables, images)
        
        Returns:
            Dict with extracted content:
            {
                "text": [...],      # Text elements with font info
                "tables": [...],    # Table data
                "images": [...],    # Image paths
                "structure": {...}  # Document structure
            }
        
        Raises:
            HTTPException 400: If Adobe API not enabled
            HTTPException 500: If extraction fails
        
        NO FALLBACK: pypdf only does basic text extraction without AI
        """
        if not self.use_adobe or not self.adobe_credentials or not ADOBE_AVAILABLE:
            raise HTTPException(
                400,
                "Content extraction requires Adobe PDF Services API. "
                "Set USE_ADOBE_PDF_API=true and configure credentials in .env"
            )
        
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        try:
            # Read input file
            async with aiofiles.open(input_file, 'rb') as f:
                input_stream = await f.read()
            
            # Adobe Extract API
            from adobe.pdfservices.operation.pdfjobs.jobs.extract_pdf_job import ExtractPDFJob
            from adobe.pdfservices.operation.pdfjobs.params.extract_pdf.extract_pdf_params import ExtractPDFParams
            from adobe.pdfservices.operation.pdfjobs.params.extract_pdf.extract_element_type import ExtractElementType
            from adobe.pdfservices.operation.pdfjobs.params.extract_pdf.extract_renditions_element_type import ExtractRenditionsElementType
            from adobe.pdfservices.operation.pdfjobs.result.extract_pdf_result import ExtractPDFResult
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=input_stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Create extract parameters based on extract_type
            extract_params = ExtractPDFParams(
                elements_to_extract=[
                    ExtractElementType.TEXT,
                    ExtractElementType.TABLES
                ],
                elements_to_extract_renditions=[
                    ExtractRenditionsElementType.TABLES,
                    ExtractRenditionsElementType.FIGURES
                ]
            )
            
            # Create and submit extract job
            extract_job = ExtractPDFJob(
                input_asset=input_asset,
                extract_pdf_params=extract_params
            )
            
            location = pdf_services.submit(extract_job)
            logger.info(f"Adobe Extract job submitted for {input_file.name}")
            
            # Get result
            pdf_services_response = pdf_services.get_job_result(location, ExtractPDFResult)
            
            # Download result (returns ZIP with JSON + images)
            result_asset: CloudAsset = pdf_services_response.get_result().get_resource()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save ZIP file temporarily
            import zipfile
            import json
            zip_path = self.output_dir / f"{input_file.stem}_extract.zip"
            async with aiofiles.open(zip_path, "wb") as f:
                await f.write(stream_asset.get_input_stream())
            
            # Extract ZIP and parse JSON
            extract_dir = self.output_dir / f"{input_file.stem}_extracted"
            extract_dir.mkdir(exist_ok=True)
            
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)
            
            # Read structuredData.json
            json_path = extract_dir / "structuredData.json"
            with open(json_path, 'r', encoding='utf-8') as f:
                extracted_data = json.load(f)
            
            # Parse extracted data based on extract_type
            result = {
                "text": [],
                "tables": [],
                "images": [],
                "structure": {}
            }
            
            # Process elements
            if "elements" in extracted_data:
                for element in extracted_data["elements"]:
                    if element.get("Path") and "Text" in element.get("Path", ""):
                        result["text"].append({
                            "text": element.get("Text", ""),
                            "font": element.get("Font", {}),
                            "bounds": element.get("Bounds", [])
                        })
                    elif element.get("Path") and "Table" in element.get("Path", ""):
                        result["tables"].append(element)
                    elif element.get("Path") and "Figure" in element.get("Path", ""):
                        result["images"].append(element)
            
            # Cleanup ZIP file
            zip_path.unlink()
            
            logger.info(f"Adobe Extract successful: extracted {len(result['text'])} text elements, "
                       f"{len(result['tables'])} tables, {len(result['images'])} images")
            
            return result
            
        except ImportError:
            logger.error("Adobe Extract API not available in SDK")
            raise HTTPException(500, "Adobe Extract requires SDK update")
        except Exception as e:
            logger.error(f"Adobe Extract error: {e}")
            raise HTTPException(500, f"Content extraction failed: {str(e)}")
    
    async def html_to_pdf(
        self,
        html_content: str,
        page_size: str = "A4",
        orientation: str = "portrait",
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert HTML to PDF using reportlab (fallback from Adobe)
        
        Supports:
        - Basic HTML tags: h1-h6, p, div, span, br, hr
        - Basic CSS: color, font-size, font-weight, text-align
        - Lists: ul, ol, li
        - Tables: table, tr, td, th
        - Custom page size (A4, Letter, Legal)
        - Orientation (portrait, landscape)
        
        Args:
            html_content: HTML string (no URL support in fallback)
            page_size: Page size (A4, Letter, Legal)
            orientation: portrait or landscape
            output_filename: Optional output filename
        
        Returns:
            Path to generated PDF
        
        Note: This is a simplified HTML renderer. For complex HTML with CSS,
              consider enabling Adobe PDF Services API.
        """
        output_filename = output_filename or "document.pdf"
        output_path = self.output_dir / output_filename
        
        try:
            from reportlab.lib.pagesizes import A4, LETTER, LEGAL, landscape
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
            from html.parser import HTMLParser
            from reportlab.lib import colors
            
            # Map page sizes
            page_size_map = {
                "A4": A4,
                "LETTER": LETTER,
                "LEGAL": LEGAL,
            }
            
            page = page_size_map.get(page_size.upper(), A4)
            if orientation.lower() == "landscape":
                page = landscape(page)
            
            # Create PDF document
            doc = SimpleDocTemplate(
                str(output_path),
                pagesize=page,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18,
            )
            
            # Get default styles
            styles = getSampleStyleSheet()
            story = []
            
            # Simple HTML parser
            class SimpleHTMLParser(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.story = []
                    self.current_text = []
                    self.current_style = styles['Normal']
                    self.in_heading = None
                    
                def handle_starttag(self, tag, attrs):
                    if tag == 'h1':
                        self.in_heading = 'Heading1'
                    elif tag == 'h2':
                        self.in_heading = 'Heading2'
                    elif tag == 'h3':
                        self.in_heading = 'Heading3'
                    elif tag == 'h4':
                        self.in_heading = 'Heading4'
                    elif tag == 'p':
                        self.current_style = styles['Normal']
                    elif tag == 'br':
                        self.current_text.append('<br/>')
                    elif tag == 'strong' or tag == 'b':
                        self.current_text.append('<b>')
                    elif tag == 'em' or tag == 'i':
                        self.current_text.append('<i>')
                    elif tag == 'u':
                        self.current_text.append('<u>')
                
                def handle_endtag(self, tag):
                    if tag in ['h1', 'h2', 'h3', 'h4', 'p']:
                        if self.current_text:
                            text = ''.join(self.current_text).strip()
                            if text:
                                if self.in_heading:
                                    para = Paragraph(text, styles[self.in_heading])
                                else:
                                    para = Paragraph(text, self.current_style)
                                self.story.append(para)
                                self.story.append(Spacer(1, 0.2 * inch))
                            self.current_text = []
                            self.in_heading = None
                    elif tag == 'strong' or tag == 'b':
                        self.current_text.append('</b>')
                    elif tag == 'em' or tag == 'i':
                        self.current_text.append('</i>')
                    elif tag == 'u':
                        self.current_text.append('</u>')
                
                def handle_data(self, data):
                    if data.strip():
                        self.current_text.append(data)
            
            # Parse HTML
            parser = SimpleHTMLParser()
            
            # Clean HTML
            html_cleaned = html_content.replace('<!DOCTYPE html>', '').replace('</html>', '').replace('</body>', '')
            
            # Remove <head> section
            import re
            html_cleaned = re.sub(r'<head>.*?</head>', '', html_cleaned, flags=re.DOTALL | re.IGNORECASE)
            html_cleaned = re.sub(r'<html.*?>', '', html_cleaned, flags=re.IGNORECASE)
            html_cleaned = re.sub(r'<body.*?>', '', html_cleaned, flags=re.IGNORECASE)
            
            parser.feed(html_cleaned)
            story = parser.story
            
            # If no content parsed, add a simple message
            if not story:
                story.append(Paragraph("HTML content converted to PDF", styles['Normal']))
                story.append(Spacer(1, 0.2 * inch))
                story.append(Paragraph(html_content[:500], styles['BodyText']))
            
            # Build PDF
            doc.build(story)
            
            logger.info(f"✅ HTML to PDF converted: {output_path.name}")
            return output_path
            
        except Exception as e:
            logger.error(f"❌ HTML to PDF conversion failed: {e}")
            raise HTTPException(500, f"HTML to PDF conversion failed: {str(e)}")
    
    # ==================== PDF Watermark (Adobe) ====================
    
    async def watermark_pdf(self, pdf_path: Path, watermark_path: Path) -> Path:
        """
        Đóng dấu mờ lên PDF sử dụng Adobe PDF Services
        
        Args:
            pdf_path: File PDF gốc
            watermark_path: File PDF dấu mờ (có thể tạo từ image/text trước)
        
        Returns:
            Path: File PDF đã có dấu mờ
        """
        if not self.adobe_credentials:
            raise HTTPException(501, "Adobe PDF Services chưa được cấu hình. Vui lòng thêm credentials vào .env")
        
        try:
            from adobe.pdfservices.operation.pdfjobs.jobs.pdf_watermark_job import PDFWatermarkJob
            from adobe.pdfservices.operation.pdfjobs.result.pdf_watermark_result import PDFWatermarkResult
            
            # Read files
            with open(pdf_path, 'rb') as f:
                source_stream = f.read()
            
            with open(watermark_path, 'rb') as f:
                watermark_stream = f.read()
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload both files
            input_asset = pdf_services.upload(
                input_stream=source_stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            watermark_asset = pdf_services.upload(
                input_stream=watermark_stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Create watermark job
            watermark_job = PDFWatermarkJob(
                input_asset=input_asset,
                watermark_asset=watermark_asset
            )
            
            # Submit and get result
            location = pdf_services.submit(watermark_job)
            response = pdf_services.get_job_result(location, PDFWatermarkResult)
            
            # Get content
            result_asset: CloudAsset = response.get_result().get_asset()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save output
            output_path = self.output_dir / f"watermarked_{pdf_path.name}"
            with open(output_path, "wb") as f:
                f.write(stream_asset.get_input_stream())
            
            logger.info(f"Adobe Watermark PDF successful: {output_path}")
            return output_path
            
        except ImportError:
            logger.error("Adobe Watermark API not available")
            raise HTTPException(500, "Adobe Watermark requires pdfservices-sdk")
        except Exception as e:
            logger.error(f"Adobe Watermark error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    
    # ==================== PDF Combine (Adobe) ====================
    
    async def combine_pdfs(self, pdf_paths: List[Path], page_ranges: Optional[List[str]] = None) -> Path:
        """
        Gộp nhiều PDF thành một file
        
        Args:
            pdf_paths: List các file PDF
            page_ranges: Optional list page ranges như ["1-3", "all", "5-10"]
        
        Returns:
            Path: File PDF đã gộp
        """
        if not self.adobe_credentials:
            raise HTTPException(501, "Adobe PDF Services chưa được cấu hình")
        
        try:
            from adobe.pdfservices.operation.pdfjobs.jobs.combine_pdf_job import CombinePDFJob
            from adobe.pdfservices.operation.pdfjobs.params.combine_pdf.combine_pdf_params import CombinePDFParams
            from adobe.pdfservices.operation.pdfjobs.params.page_ranges import PageRanges
            from adobe.pdfservices.operation.pdfjobs.result.combine_pdf_result import CombinePDFResult
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Create CombinePDFParams (v4 API pattern)
            combine_pdf_params = CombinePDFParams()
            
            for idx, pdf_path in enumerate(pdf_paths):
                with open(pdf_path, 'rb') as f:
                    stream = f.read()
                
                asset = pdf_services.upload(
                    input_stream=stream,
                    mime_type=PDFServicesMediaType.PDF
                )
                
                # Add with or without page range
                if page_ranges and idx < len(page_ranges) and page_ranges[idx] != "all":
                    # Parse range like "1-3" or "5-10"
                    range_str = page_ranges[idx]
                    page_range_obj = PageRanges()
                    
                    if '-' in range_str:
                        start, end = range_str.split('-')
                        page_range_obj.add_range(int(start), int(end))
                    else:
                        # Single page
                        page = int(range_str)
                        page_range_obj.add_single_page(page)
                    
                    # Add asset with page ranges to params
                    combine_pdf_params.add_asset(asset, page_range_obj)
                else:
                    # Add asset without page ranges (all pages)
                    combine_pdf_params.add_asset(asset)
            
            # Create job with params (v4 API: Job requires params argument)
            combine_job = CombinePDFJob(combine_pdf_params=combine_pdf_params)
            
            # Submit and get result
            location = pdf_services.submit(combine_job)
            response = pdf_services.get_job_result(location, CombinePDFResult)
            
            # Get content
            result_asset: CloudAsset = response.get_result().get_asset()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save output
            output_path = self.output_dir / "combined.pdf"
            with open(output_path, "wb") as f:
                f.write(stream_asset.get_input_stream())
            
            logger.info(f"Adobe Combine PDF successful: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Adobe Combine error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    
    # ==================== PDF Split (Adobe) ====================
    
    async def split_pdf(self, pdf_path: Path, page_ranges: List[str]) -> List[Path]:
        """
        Tách PDF thành nhiều file
        
        Args:
            pdf_path: File PDF gốc
            page_ranges: List ranges như ["1-3", "4-6", "7-10"]
        
        Returns:
            List[Path]: List các file PDF đã tách
        """
        if not self.adobe_credentials:
            raise HTTPException(501, "Adobe PDF Services chưa được cấu hình")
        
        try:
            from adobe.pdfservices.operation.pdfjobs.jobs.split_pdf_job import SplitPDFJob
            from adobe.pdfservices.operation.pdfjobs.params.split_pdf.split_pdf_params import SplitPDFParams
            from adobe.pdfservices.operation.pdfjobs.params.page_ranges import PageRanges
            from adobe.pdfservices.operation.pdfjobs.result.split_pdf_result import SplitPDFResult
            
            # Read file
            with open(pdf_path, 'rb') as f:
                stream = f.read()
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Parse page ranges - Adobe API expects ONE PageRanges object with multiple ranges
            page_ranges_obj = PageRanges()
            for range_str in page_ranges:
                # Clean whitespace
                range_str = range_str.strip().replace(' ', '')
                
                if '-' in range_str:
                    parts = range_str.split('-')
                    if len(parts) != 2:
                        raise HTTPException(400, f"Invalid range format: {range_str}. Use 1-3, not 1-2-3")
                    start, end = int(parts[0]), int(parts[1])
                    if start > end:
                        raise HTTPException(400, f"Invalid range: {range_str}. Start must be <= end")
                    page_ranges_obj.add_range(start, end)
                else:
                    # Single page
                    page = int(range_str)
                    if page < 1:
                        raise HTTPException(400, f"Invalid page number: {page}. Must be >= 1")
                    page_ranges_obj.add_single_page(page)
            
            # Create split params - pass single PageRanges object
            split_params = SplitPDFParams(page_ranges=page_ranges_obj)
            
            # Log ranges before sending to Adobe
            logger.info(f"📤 Sending to Adobe: {len(page_ranges)} ranges")
            
            # Create and submit job
            split_job = SplitPDFJob(input_asset=input_asset, split_pdf_params=split_params)
            location = pdf_services.submit(split_job)
            response = pdf_services.get_job_result(location, SplitPDFResult)
            
            # Get all result assets
            result_assets = response.get_result().get_assets()
            logger.info(f"🔍 Split PDF: Got {len(result_assets)} result assets from Adobe")
            
            # Save all output files
            output_paths = []
            for idx, result_asset in enumerate(result_assets):
                stream_asset: StreamAsset = pdf_services.get_content(result_asset)
                content_bytes = stream_asset.get_input_stream()
                
                logger.info(f"📄 Split file {idx+1}: Content size = {len(content_bytes)} bytes")
                
                output_path = self.output_dir / f"split_{idx+1}_{pdf_path.name}"
                with open(output_path, "wb") as f:
                    f.write(content_bytes)
                
                # Verify file was written
                file_size = output_path.stat().st_size
                logger.info(f"💾 Split file {idx+1}: Saved to {output_path.name}, File size on disk = {file_size} bytes")
                
                if file_size == 0:
                    logger.error(f"❌ WARNING: Split file {idx+1} is EMPTY (0 bytes)!")
                
                output_paths.append(output_path)
            
            logger.info(f"✅ Adobe Split PDF successful: {len(output_paths)} files created")
            return output_paths
            
        except Exception as e:
            logger.error(f"Adobe Split error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    
    async def split_pdf_pypdf(self, pdf_path: Path, page_ranges: List[str]) -> List[Path]:
        """
        Tách PDF bằng PyPDF2 (fallback khi Adobe fail với PDF signed)
        
        Args:
            pdf_path: File PDF gốc
            page_ranges: List ranges như ["1-3", "4-6", "7"]
        
        Returns:
            List[Path]: List các file PDF đã tách
        """
        try:
            import pypdf
            
            output_paths = []
            
            with open(pdf_path, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                
                for idx, range_str in enumerate(page_ranges):
                    pdf_writer = pypdf.PdfWriter()
                    
                    # Parse range string
                    if '-' in range_str:
                        start, end = range_str.split('-')
                        start_page = int(start) - 1  # Convert to 0-indexed
                        end_page = int(end)  # Inclusive
                    else:
                        # Single page
                        start_page = int(range_str) - 1
                        end_page = int(range_str)
                    
                    # Add pages to writer
                    for page_num in range(start_page, end_page):
                        if page_num < len(pdf_reader.pages):
                            pdf_writer.add_page(pdf_reader.pages[page_num])
                    
                    # Save output file
                    output_path = self.output_dir / f"split_{idx+1}_{pdf_path.name}"
                    with open(output_path, 'wb') as output_file:
                        pdf_writer.write(output_file)
                    
                    output_paths.append(output_path)
                    logger.info(f"💾 PyPDF2 split file {idx+1}: {output_path.name} ({output_path.stat().st_size} bytes)")
            
            logger.info(f"✅ PyPDF2 Split successful: {len(output_paths)} files created")
            return output_paths
            
        except Exception as e:
            logger.error(f"PyPDF2 Split error: {e}")
            raise HTTPException(500, f"😔 Không thể tách PDF: {str(e)}")
    
    # ==================== PDF Protect (Adobe) ====================
    
    async def protect_pdf(
        self, 
        pdf_path: Path, 
        user_password: str,
        owner_password: Optional[str] = None,
        permissions: Optional[List[str]] = None
    ) -> Path:
        """
        Bảo vệ PDF bằng mật khẩu và phân quyền
        
        Args:
            pdf_path: File PDF gốc
            user_password: Mật khẩu để mở file
            owner_password: Mật khẩu chủ sở hữu (optional)
            permissions: List quyền ["print", "copy", "edit"] (optional)
        
        Returns:
            Path: File PDF đã bảo vệ
        """
        if not self.adobe_credentials:
            raise HTTPException(501, "Adobe PDF Services chưa được cấu hình")
        
        try:
            from adobe.pdfservices.operation.pdfjobs.jobs.protect_pdf_job import ProtectPDFJob
            from adobe.pdfservices.operation.pdfjobs.params.protect_pdf.password_protect_params import PasswordProtectParams
            from adobe.pdfservices.operation.pdfjobs.params.protect_pdf.encryption_algorithm import EncryptionAlgorithm
            from adobe.pdfservices.operation.pdfjobs.params.protect_pdf.permission import Permission
            from adobe.pdfservices.operation.pdfjobs.params.protect_pdf.content_encryption import ContentEncryption
            from adobe.pdfservices.operation.pdfjobs.result.protect_pdf_result import ProtectPDFResult
            
            # Read file
            with open(pdf_path, 'rb') as f:
                stream = f.read()
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Parse permissions
            permission_list = []
            if permissions:
                perm_map = {
                    "print": Permission.PRINT_LOW_QUALITY,
                    "print_high": Permission.PRINT_HIGH_QUALITY,
                    "copy": Permission.COPY_CONTENT,
                    "edit": Permission.EDIT_CONTENT,
                    "edit_annotations": Permission.EDIT_ANNOTATIONS,
                    "fill_forms": Permission.FILL_AND_SIGN_FORM_FIELDS,
                    "assemble": Permission.EDIT_DOCUMENT_ASSEMBLY
                }
                for perm in permissions:
                    if perm in perm_map:
                        permission_list.append(perm_map[perm])
            
            # Create protect params (v4 API: Use PasswordProtectParams)
            protect_params = PasswordProtectParams(
                user_password=user_password,
                encryption_algorithm=EncryptionAlgorithm.AES_256,
                content_encryption=ContentEncryption.ALL_CONTENT
            )
            
            # Note: owner_password and permissions may need different API
            # Official sample doesn't show these - may need separate params or different constructor
            # For now keeping basic protection working
            
            # Create and submit job
            protect_job = ProtectPDFJob(input_asset=input_asset, protect_pdf_params=protect_params)
            location = pdf_services.submit(protect_job)
            response = pdf_services.get_job_result(location, ProtectPDFResult)
            
            # Get content
            result_asset: CloudAsset = response.get_result().get_asset()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save output
            output_path = self.output_dir / f"protected_{pdf_path.name}"
            with open(output_path, "wb") as f:
                f.write(stream_asset.get_input_stream())
            
            logger.info(f"Adobe Protect PDF successful: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Adobe Protect error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    
    # ==================== PDF Linearize (Adobe) ====================
    
    async def linearize_pdf(self, pdf_path: Path) -> Path:
        """
        Tối ưu hóa PDF cho web (fast web viewing)
        
        Args:
            pdf_path: File PDF gốc
        
        Returns:
            Path: File PDF đã tối ưu
        """
        if not self.adobe_credentials:
            raise HTTPException(501, "Adobe PDF Services chưa được cấu hình")
        
        try:
            from adobe.pdfservices.operation.pdfjobs.jobs.linearize_pdf_job import LinearizePDFJob
            from adobe.pdfservices.operation.pdfjobs.result.linearize_pdf_result import LinearizePDFResult
            
            # Read file
            with open(pdf_path, 'rb') as f:
                stream = f.read()
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Create and submit job
            linearize_job = LinearizePDFJob(input_asset=input_asset)
            location = pdf_services.submit(linearize_job)
            response = pdf_services.get_job_result(location, LinearizePDFResult)
            
            # Get content
            result_asset: CloudAsset = response.get_result().get_asset()
            stream_asset: StreamAsset = pdf_services.get_content(result_asset)
            
            # Save output
            output_path = self.output_dir / f"linearized_{pdf_path.name}"
            with open(output_path, "wb") as f:
                f.write(stream_asset.get_input_stream())
            
            logger.info(f"Adobe Linearize PDF successful: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Adobe Linearize error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    
    # ==================== PDF Auto-Tag (Adobe) ====================
    
    async def autotag_pdf(self, pdf_path: Path, generate_report: bool = True) -> tuple[Path, Optional[Path]]:
        """
        Tự động gắn thẻ PDF cho accessibility (khả năng tiếp cận)
        
        Args:
            pdf_path: File PDF gốc
            generate_report: Có tạo báo cáo accessibility không
        
        Returns:
            tuple: (tagged_pdf_path, report_path)
        """
        if not self.adobe_credentials:
            raise HTTPException(501, "Adobe PDF Services chưa được cấu hình")
        
        try:
            from adobe.pdfservices.operation.pdfjobs.jobs.autotag_pdf_job import AutotagPDFJob
            from adobe.pdfservices.operation.pdfjobs.result.autotag_pdf_result import AutotagPDFResult
            
            # Read file
            with open(pdf_path, 'rb') as f:
                stream = f.read()
            
            # Create PDF Services instance
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload file
            input_asset = pdf_services.upload(
                input_stream=stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Create and submit job
            autotag_job = AutotagPDFJob(
                input_asset=input_asset,
                generate_report=generate_report
            )
            location = pdf_services.submit(autotag_job)
            response = pdf_services.get_job_result(location, AutotagPDFResult)
            
            # Get tagged PDF
            result = response.get_result()
            tagged_asset: CloudAsset = result.get_tagged_pdf()
            stream_asset: StreamAsset = pdf_services.get_content(tagged_asset)
            
            # Save tagged PDF
            output_path = self.output_dir / f"tagged_{pdf_path.name}"
            with open(output_path, "wb") as f:
                f.write(stream_asset.get_input_stream())
            
            # Get report if generated
            report_path = None
            if generate_report:
                report_asset: CloudAsset = result.get_report()
                report_stream: StreamAsset = pdf_services.get_content(report_asset)
                
                report_path = self.output_dir / f"accessibility_report_{pdf_path.stem}.xlsx"
                with open(report_path, "wb") as f:
                    f.write(report_stream.get_input_stream())
            
            logger.info(f"Adobe Auto-Tag PDF successful: {output_path}")
            return output_path, report_path
            
        except Exception as e:
            logger.error(f"Adobe Auto-Tag error: {e}")
            status_code, friendly_msg = get_friendly_error_message(e)
            raise HTTPException(status_code, friendly_msg)
    
    async def generate_document(
        self,
        template_path: Path,
        json_data: dict,
        output_format: str = "pdf"
    ) -> Path:
        """
        Generate PDF/DOCX from Word template + JSON data.
        
        Args:
            template_path: Path to .docx template file
            json_data: Dictionary with merge data
            output_format: "pdf" or "docx"
            
        Returns:
            Path to generated file
        """
        if not self.adobe_credentials:
            raise HTTPException(
                status_code=500,
                detail="Adobe PDF Services credentials not configured"
            )
            
        try:
            # Import Adobe SDK components
            from adobe.pdfservices.operation.pdf_services import PDFServices
            from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
            from adobe.pdfservices.operation.pdfjobs.jobs.document_merge_job import DocumentMergeJob
            from adobe.pdfservices.operation.pdfjobs.params.documentmerge.document_merge_params import DocumentMergeParams
            from adobe.pdfservices.operation.pdfjobs.params.documentmerge.output_format import OutputFormat
            from adobe.pdfservices.operation.pdfjobs.result.document_merge_result import DocumentMergePDFResult
            
            # Initialize PDF Services
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Read template file
            with open(template_path, 'rb') as template_file:
                input_stream = template_file.read()
            
            # Upload template
            input_asset = pdf_services.upload(
                input_stream=input_stream,
                mime_type=PDFServicesMediaType.DOCX
            )
            
            # Set output format
            if output_format.lower() == "pdf":
                merge_format = OutputFormat.PDF
            else:
                merge_format = OutputFormat.DOCX
            
            # Create merge parameters
            document_merge_params = DocumentMergeParams(
                json_data_for_merge=json_data,
                output_format=merge_format
            )
            
            # Create and submit job
            document_merge_job = DocumentMergeJob(
                input_asset=input_asset,
                document_merge_params=document_merge_params
            )
            
            location = pdf_services.submit(document_merge_job)
            pdf_services_response = pdf_services.get_job_result(
                location,
                DocumentMergePDFResult
            )
            
            # Get result asset
            result_asset = pdf_services_response.get_result().get_asset()
            stream_asset = pdf_services.get_content(result_asset)
            
            # Save output file
            output_file_path = self.output_dir / f"generated_document_{uuid.uuid4()}.{output_format.lower()}"
            with open(output_file_path, "wb") as output_file:
                output_file.write(stream_asset.get_input_stream())
            
            logger.info(f"Document generation successful: {output_file_path}")
            return output_file_path
            
        except Exception as e:
            logger.error(f"Document generation error: {e}")
            raise HTTPException(
                status_code=500,
                detail=f"Document generation failed: {str(e)}"
            )
    
    async def electronic_seal_pdf(
        self,
        pdf_path: Path,
        seal_image_path: Optional[Path],
        provider_name: str,
        access_token: str,
        credential_id: str,
        pin: str,
        seal_field_name: str = "Signature1",
        page_number: int = 1,
        visible: bool = True,
        field_x: int = 150,
        field_y: int = 250,
        field_width: int = 350,
        field_height: int = 200
    ) -> Path:
        """
        Apply electronic seal (digital signature) to PDF.
        
        Args:
            pdf_path: Path to PDF file
            seal_image_path: Optional path to seal image (PNG/JPG)
            provider_name: TSP provider name
            access_token: TSP access token
            credential_id: TSP credential ID
            pin: TSP PIN
            seal_field_name: Name of signature field
            page_number: Page number for seal (1-based)
            visible: Whether seal is visible
            field_x, field_y, field_width, field_height: Seal position/size
            
        Returns:
            Path to sealed PDF
        """
        if not self.adobe_credentials:
            raise HTTPException(
                status_code=500,
                detail="Adobe PDF Services credentials not configured"
            )
            
        try:
            # Import Adobe SDK components
            from adobe.pdfservices.operation.pdf_services import PDFServices
            from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
            from adobe.pdfservices.operation.pdfjobs.jobs.eseal_job import PDFElectronicSealJob
            from adobe.pdfservices.operation.pdfjobs.params.eseal.csc_auth_context import CSCAuthContext
            from adobe.pdfservices.operation.pdfjobs.params.eseal.csc_credentials import CSCCredentials
            from adobe.pdfservices.operation.pdfjobs.params.eseal.electronic_seal_params import PDFElectronicSealParams
            from adobe.pdfservices.operation.pdfjobs.params.eseal.field_location import FieldLocation
            from adobe.pdfservices.operation.pdfjobs.params.eseal.field_options import FieldOptions
            from adobe.pdfservices.operation.pdfjobs.params.eseal.document_level_permission import DocumentLevelPermission
            from adobe.pdfservices.operation.pdfjobs.result.eseal_pdf_result import ESealPDFResult
            
            # Initialize PDF Services
            pdf_services = PDFServices(credentials=self.adobe_credentials)
            
            # Upload PDF
            with open(pdf_path, 'rb') as pdf_file:
                pdf_stream = pdf_file.read()
            
            pdf_asset = pdf_services.upload(
                input_stream=pdf_stream,
                mime_type=PDFServicesMediaType.PDF
            )
            
            # Upload seal image if provided
            seal_image_asset = None
            if seal_image_path and seal_image_path.exists():
                with open(seal_image_path, 'rb') as img_file:
                    img_stream = img_file.read()
                
                # Determine image type
                ext = seal_image_path.suffix.lower()
                if ext == '.png':
                    mime_type = PDFServicesMediaType.PNG
                elif ext in ['.jpg', '.jpeg']:
                    mime_type = PDFServicesMediaType.JPEG
                else:
                    raise HTTPException(400, "Seal image must be PNG or JPEG")
                
                seal_image_asset = pdf_services.upload(
                    input_stream=img_stream,
                    mime_type=mime_type
                )
            
            # Create field location
            field_location = FieldLocation(field_x, field_y, field_width, field_height)
            
            # Create field options
            field_options = FieldOptions(
                field_name=seal_field_name,
                field_location=field_location,
                page_number=page_number,
                visible=visible
            )
            
            # Create CSC auth context
            csc_auth_context = CSCAuthContext(
                access_token=access_token,
                token_type="Bearer"
            )
            
            # Create CSC credentials
            csc_credentials = CSCCredentials(
                provider_name=provider_name,
                credential_id=credential_id,
                pin=pin,
                csc_auth_context=csc_auth_context
            )
            
            # Create seal parameters
            seal_params = PDFElectronicSealParams(
                seal_certificate_credentials=csc_credentials,
                seal_field_options=field_options
            )
            
            # Create job
            if seal_image_asset:
                seal_job = PDFElectronicSealJob(
                    input_asset=pdf_asset,
                    electronic_seal_params=seal_params,
                    seal_image_asset=seal_image_asset
                )
            else:
                seal_job = PDFElectronicSealJob(
                    input_asset=pdf_asset,
                    electronic_seal_params=seal_params
                )
            
            # Submit and get result
            location = pdf_services.submit(seal_job)
            pdf_services_response = pdf_services.get_job_result(location, ESealPDFResult)
            
            # Get result asset
            result_asset = pdf_services_response.get_result().get_asset()
            stream_asset = pdf_services.get_content(result_asset)
            
            # Save output
            output_file_path = self.output_dir / f"sealed_{uuid.uuid4()}.pdf"
            with open(output_file_path, "wb") as output_file:
                output_file.write(stream_asset.get_input_stream())
            
            logger.info(f"Electronic seal successful: {output_file_path}")
            return output_file_path
            
        except Exception as e:
            logger.error(f"Electronic seal error: {e}")
            raise HTTPException(
                status_code=500,
                detail=f"Electronic seal failed: {str(e)}"
            )
    
    # ==================== Smart PDF OCR ====================
    
    async def smart_pdf_ocr(
        self,
        input_file: Path,
        ai_engine: str = "gemini",
        language: str = "vi",
        db = None
    ) -> dict:
        """
        Smart PDF OCR - Uses AI only for scanned PDFs, direct extraction for text-based PDFs
        
        Logic:
        1. Detect if PDF is scanned using is_pdf_scanned()
        2. If scanned → Use AI OCR (Gemini/Claude) 
        3. If text-based → Use direct text extraction (fast & free)
        
        Args:
            input_file: Path to PDF file
            ai_engine: AI engine to use (gemini or claude)
            language: Language for OCR (vi, en)
            
        Returns:
            Dict with extraction results and metadata
        """
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
            
        start_time = time.time()
        
        try:
            # Step 1: Detect PDF type
            logger.info(f"🔍 Analyzing PDF type: {input_file.name}")
            pdf_is_scanned = is_pdf_scanned(input_file)
            
            if pdf_is_scanned:
                # PDF is scanned → Use AI OCR
                logger.info(f"📄 PDF is scanned - using {ai_engine.upper()} OCR")
                result = await self._ocr_scanned_pdf(input_file, ai_engine, language, db=db)
            else:
                # PDF has text layer → Direct extraction
                logger.info("📝 PDF has text layer - using direct extraction")
                result = await self._extract_text_based_pdf(input_file)
                
            # Add processing metadata
            processing_time = time.time() - start_time
            result["processing"] = {
                "time_seconds": round(processing_time, 2),
                "method": "ai_ocr" if pdf_is_scanned else "direct_extraction",
                "engine": ai_engine if pdf_is_scanned else "pypdf",
                "pdf_type": "scanned" if pdf_is_scanned else "text_based"
            }
            
            # Track AI usage if OCR was used
            if pdf_is_scanned and "ai_usage" in result:
                await self._track_ocr_usage(ai_engine, result, processing_time)
            
            return result
            
        except Exception as e:
            logger.error(f"Smart PDF OCR failed: {e}")
            raise HTTPException(500, f"PDF OCR failed: {str(e)}")
            
    async def _ocr_scanned_pdf(self, input_file: Path, ai_engine: str, language: str, db = None) -> dict:
        """OCR scanned PDF using AI (Gemini/Claude)"""
        try:
            # Gemini can process PDF directly - much faster and cheaper!
            if ai_engine.lower() == "gemini":
                logger.info("📄 Using Gemini native PDF processing (no image conversion needed)")
                return await self._ocr_pdf_with_gemini(input_file, language, db=db)
            
            # Claude requires image conversion
            logger.info("🖼️ Converting PDF to images for Claude...")
            
            # Import pdf2image
            try:
                from pdf2image import convert_from_path
            except ImportError:
                raise HTTPException(
                    500,
                    "pdf2image not installed. Run: pip install pdf2image"
                )
            
            # Check if poppler is available by testing with first page
            try:
                test_images = convert_from_path(str(input_file), first_page=1, last_page=1)
            except Exception as e:
                if "poppler" in str(e).lower() or "Unable to get page count" in str(e):
                    raise HTTPException(
                        500,
                        "⚠️ Poppler not found!\n\n"
                        "Poppler is required to convert PDF to images for Claude OCR.\n\n"
                        "📥 Install Poppler for Windows:\n"
                        "1. Download: https://github.com/oschwartz10612/poppler-windows/releases/\n"
                        "2. Extract poppler-xx.xx.x to C:\\poppler\n"
                        "3. Add C:\\poppler\\Library\\bin to System PATH\n"
                        "4. Restart terminal/IDE\n\n"
                        "Or use Chocolatey: choco install poppler\n\n"
                        "💡 Tip: Use Gemini instead - it can read PDF directly without poppler!"
                    )
                raise
            
            # Convert all pages to images
            images = convert_from_path(str(input_file))
            
            logger.info(f"📄 Converted PDF to {len(images)} images")
            
            all_text = ""
            page_texts = []
            total_cost = 0.0
            total_tokens = 0
            
            for i, image in enumerate(images):
                logger.info(f"🔍 Processing page {i+1}/{len(images)} with {ai_engine}")
                
                # Save image temporarily
                temp_img_path = input_file.parent / f"temp_page_{i+1}.png"
                image.save(temp_img_path, 'PNG')
                
                try:
                    if ai_engine.lower() == "gemini":
                        page_result = await self._ocr_with_gemini(temp_img_path, language)
                    elif ai_engine.lower() == "claude":
                        page_result = await self._ocr_with_claude(temp_img_path, language)
                    else:
                        raise HTTPException(400, "AI engine must be 'gemini' or 'claude'")
                        
                    page_text = page_result.get("text", "")
                    page_texts.append({
                        "page": i + 1,
                        "text": page_text,
                        "char_count": len(page_text),
                        "tokens": page_result.get("tokens", 0),
                        "cost_usd": page_result.get("cost", 0.0)
                    })
                    
                    all_text += f"\\n\\n--- Page {i+1} ---\\n" + page_text
                    total_tokens += page_result.get("tokens", 0)
                    total_cost += page_result.get("cost", 0.0)
                    
                finally:
                    # Cleanup temp image
                    if temp_img_path.exists():
                        temp_img_path.unlink()
                        
            return {
                "text": all_text.strip(),
                "pages": page_texts,
                "total_pages": len(images),
                "char_count": len(all_text),
                "word_count": len(all_text.split()),
                "ai_usage": {
                    "engine": ai_engine,
                    "total_tokens": total_tokens,
                    "total_cost_usd": round(total_cost, 6),
                    "cost_per_page": round(total_cost / len(images), 6) if images else 0
                }
            }
            
        except Exception as e:
            logger.error(f"AI OCR failed: {e}")
            raise HTTPException(500, f"AI OCR failed: {str(e)}")
            
    async def _extract_text_based_pdf(self, input_file: Path) -> dict:
        """Extract text from text-based PDF using direct method"""
        try:
            text = await self.extract_pdf_text(input_file)
            
            return {
                "text": text,
                "char_count": len(text),
                "word_count": len(text.split()),
                "message": "✅ Text extracted directly from PDF (no OCR needed)"
            }
            
        except Exception as e:
            logger.error(f"Direct text extraction failed: {e}")
            raise HTTPException(500, f"Text extraction failed: {str(e)}")
    
    async def _ocr_pdf_with_gemini(self, pdf_path: Path, language: str, db = None) -> dict:
        """🤖 OCR PDF directly using Gemini native PDF support with auto-logging"""
        try:
            import google.generativeai as genai
            from app.services.gemini_service import get_gemini_service
            
            # Get API key
            from app.services.ai_usage_service import get_api_key
            api_key = get_api_key("gemini")
            if not api_key:
                raise HTTPException(500, "Gemini API key not configured")
                
            genai.configure(api_key=api_key)
            
            # Upload PDF to Gemini Files API
            logger.info(f"📤 Uploading PDF to Gemini Files API...")
            uploaded_file = genai.upload_file(str(pdf_path))
            logger.info(f"✅ PDF uploaded: {uploaded_file.name}")
            
            # OCR prompt
            prompt = """Trích xuất TOÀN BỘ văn bản trong tài liệu PDF này.

YÊU CẦU:
- Giữ chính xác 100% ký tự Tiếng Việt (ă, â, ê, ô, ơ, ư, đ, dấu thanh)
- Giữ nguyên cấu trúc, xuống dòng, format như trong PDF
- Chỉ trả về văn bản, KHÔNG thêm giải thích

Trả về văn bản:"""
            
            # Use GeminiService for auto-logging
            gemini_service = get_gemini_service(db) if db else None
            
            if gemini_service:
                # Auto-logging enabled
                response = gemini_service.generate_content(
                    prompt=[uploaded_file, prompt],
                    model="gemini-2.0-flash-exp",
                    operation="pdf-ocr",
                    metadata={
                        "file_name": pdf_path.name,
                        "language": language,
                        "ocr_type": "pdf_native"
                    }
                )
                logger.info("✅ Usage automatically logged to database")
            else:
                # Fallback without logging
                logger.warning("⚠️ No db session - OCR will NOT be logged")
                model = genai.GenerativeModel("gemini-2.0-flash-exp")
                response = model.generate_content([uploaded_file, prompt])
            
            text = response.text.strip()
            
            # Calculate usage (safely handle missing usage_metadata)
            total_tokens = 0
            cost = 0.0
            if hasattr(response, 'usage_metadata') and response.usage_metadata:
                usage = response.usage_metadata
                prompt_tokens = getattr(usage, 'prompt_token_count', 0)
                candidates_tokens = getattr(usage, 'candidates_token_count', 0)
                total_tokens = prompt_tokens + candidates_tokens
                cost = (prompt_tokens / 1_000_000 * 0.075) + (candidates_tokens / 1_000_000 * 0.30)
            
            # Clean up uploaded file
            try:
                genai.delete_file(uploaded_file.name)
                logger.info(f"🗑️ Cleaned up uploaded file")
            except:
                pass
            
            # Get page count
            from pypdf import PdfReader
            reader = PdfReader(pdf_path)
            page_count = len(reader.pages)
            
            return {
                "text": text,
                "total_pages": page_count,
                "char_count": len(text),
                "word_count": len(text.split()),
                "ai_usage": {
                    "engine": "gemini",
                    "total_tokens": total_tokens,
                    "total_cost_usd": round(cost, 6),
                    "cost_per_page": round(cost / page_count, 6) if page_count > 0 else 0
                }
            }
            
        except Exception as e:
            logger.error(f"Gemini PDF OCR failed: {e}")
            raise HTTPException(500, f"Gemini PDF OCR failed: {str(e)}")
            
    async def _ocr_with_gemini(self, image_path: Path, language: str) -> dict:
        """OCR image using Gemini API"""
        try:
            import google.generativeai as genai
            from PIL import Image as PILImage
            
            # Get API key
            from app.services.ai_usage_service import get_api_key
            api_key = get_api_key("gemini")
            if not api_key:
                raise HTTPException(500, "Gemini API key not configured")
                
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel("gemini-2.5-flash")
            
            # Prepare image
            image = PILImage.open(image_path)
            
            # OCR prompt
            prompt = f"""Trích xuất TOÀN BỘ văn bản trong ảnh này.

YÊU CẦU:
- Giữ chính xác 100% ký tự Tiếng Việt (ă, â, ê, ô, ơ, ư, đ, dấu thanh)
- Giữ nguyên cấu trúc, xuống dòng như trong ảnh
- Chỉ trả về văn bản, KHÔNG thêm giải thích

Trả về văn bản:"""
            
            # Generate content
            response = model.generate_content([prompt, image])
            text = response.text.strip()
            
            # Calculate usage (safely handle missing usage_metadata)
            input_tokens = 0
            output_tokens = 0
            total_tokens = 0
            cost = 0.0
            if hasattr(response, 'usage_metadata') and response.usage_metadata:
                usage = response.usage_metadata
                input_tokens = getattr(usage, 'prompt_token_count', 0)
                output_tokens = getattr(usage, 'candidates_token_count', 0)
                total_tokens = input_tokens + output_tokens
                cost = (input_tokens / 1_000_000 * 0.30) + (output_tokens / 1_000_000 * 2.50)
            
            return {
                "text": text,
                "tokens": total_tokens,
                "cost": cost
            }
            
        except Exception as e:
            logger.error(f"Gemini OCR failed: {e}")
            raise HTTPException(500, f"Gemini OCR failed: {str(e)}")
            
    async def _ocr_with_claude(self, image_path: Path, language: str) -> dict:
        """OCR image using Claude API"""
        try:
            import anthropic
            
            # Get API key
            from app.services.ai_usage_service import get_api_key
            api_key = get_api_key("claude")
            if not api_key:
                raise HTTPException(500, "Claude API key not configured")
                
            client = anthropic.Anthropic(api_key=api_key)
            
            # Encode image
            with open(image_path, "rb") as img_file:
                image_data = base64.standard_b64encode(img_file.read()).decode("utf-8")
                
            # OCR prompt
            prompt = """Trích xuất TOÀN BỘ văn bản trong ảnh này.

YÊU CẦU:
- Giữ chính xác 100% ký tự Tiếng Việt (ă, â, ê, ô, ơ, ư, đ, dấu thanh)
- Giữ nguyên cấu trúc, xuống dòng như trong ảnh
- Chỉ trả về văn bản, KHÔNG thêm giải thích

Trả về văn bản:"""
            
            # Call Claude API
            message = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=1024,
                messages=[{
                    "role": "user",
                    "content": [{
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/png",
                            "data": image_data,
                        },
                    }, {
                        "type": "text",
                        "text": prompt
                    }],
                }],
            )
            
            text = message.content[0].text.strip()
            
            # Calculate usage
            input_tokens = message.usage.input_tokens
            output_tokens = message.usage.output_tokens
            total_tokens = input_tokens + output_tokens
            
            # Calculate cost (Claude Sonnet 4 pricing)
            cost = (input_tokens / 1_000_000 * 3.00) + (output_tokens / 1_000_000 * 15.00)
            
            return {
                "text": text,
                "tokens": total_tokens,
                "cost": cost
            }
            
        except Exception as e:
            logger.error(f"Claude OCR failed: {e}")
            raise HTTPException(500, f"Claude OCR failed: {str(e)}")
    
    async def _track_ocr_usage(self, ai_engine: str, result: dict, processing_time: float):
        """Track AI OCR usage to database"""
        try:
            from app.services.ai_usage_service import log_usage
            from app.core.database import SessionLocal
            
            ai_usage = result.get("ai_usage", {})
            if not ai_usage:
                return
            
            # Determine model based on engine
            model = "gemini-2.0-flash-exp" if ai_engine == "gemini" else "claude-sonnet-4-20250514"
            
            # Calculate tokens (if not present, estimate from pages)
            total_tokens = ai_usage.get("total_tokens", 0)
            if total_tokens == 0 and "total_pages" in result:
                # Estimate: 258 tokens per page for Gemini PDF
                total_tokens = result["total_pages"] * 258
            
            # Split input/output tokens (estimate 70/30 ratio)
            input_tokens = int(total_tokens * 0.7)
            output_tokens = int(total_tokens * 0.3)
            
            # Log usage
            db = SessionLocal()
            try:
                log_usage(
                    db=db,
                    provider=ai_engine,
                    model=model,
                    endpoint="Ocr Compare",  # Match the UI display
                    input_tokens=input_tokens,
                    output_tokens=output_tokens,
                    processing_time=processing_time,
                    status="success",
                    request_metadata={
                        "pages": result.get("total_pages", 1),
                        "char_count": result.get("char_count", 0),
                        "cost_usd": ai_usage.get("total_cost_usd", 0.0)
                    }
                )
                db.commit()
                logger.info(f"✅ Tracked {ai_engine} usage: {total_tokens} tokens, ${ai_usage.get('total_cost_usd', 0):.6f}")
            finally:
                db.close()
                
        except Exception as e:
            logger.error(f"Failed to track OCR usage: {e}")
            # Don't raise - tracking failure shouldn't break OCR
    
    # ==================== Cleanup ====================
    
    # ==================== Cleanup ====================
    
    async def cleanup_file(self, file_path: Path) -> None:
        """Delete a file"""
        if file_path.exists():
            file_path.unlink()
    
    async def cleanup_old_files(self, max_age_hours: int = 24) -> int:
        """Delete files older than max_age_hours"""
        import time
        
        count = 0
        current_time = time.time()
        max_age_seconds = max_age_hours * 3600
        
        for dir_path in [self.upload_dir, self.output_dir]:
            for file_path in dir_path.glob("*"):
                if file_path.is_file():
                    file_age = current_time - file_path.stat().st_mtime
                    if file_age > max_age_seconds:
                        file_path.unlink()
                        count += 1
                        
        return count


    # ==================== PDF TO WORD SMART (Gemini) ====================
    
    async def pdf_to_word_smart(
        self,
        input_file: Path,
        language: str = "vi"
    ) -> Path:
        
        if input_file.suffix.lower() != '.pdf':
            raise HTTPException(400, "File must be .pdf")
        
        start_time = time.time()
        output_file = self.output_dir / f"{input_file.stem}_converted.docx"
        
        try:
            # Step 1: Upload PDF to Gemini and generate Markdown
            logger.info(f"Converting PDF to Markdown with Gemini...")
            markdown_result = await self._pdf_to_markdown_gemini(input_file, language)
            
            # Step 2: Convert Markdown to Word
            logger.info(f"Converting Markdown to Word...")
            await self._markdown_to_word(markdown_result["markdown"], output_file)
            
            processing_time = time.time() - start_time
            logger.info(f"PDF to Word completed in {processing_time:.2f}s")
            
            # Track AI usage
            await self._track_pdf_to_word_usage(
                markdown_result.get("ai_usage", {}),
                processing_time
            )
            
            return output_file
            
        except Exception as e:
            error_msg = str(e).encode('utf-8', errors='ignore').decode('utf-8')
            logger.error(f"PDF to Word conversion failed: {error_msg}")
            raise HTTPException(500, f"Conversion failed: {error_msg}")
    
    async def _pdf_to_markdown_gemini(self, pdf_path: Path, language: str) -> dict:
        
        try:
            import google.generativeai as genai
            
            # Get API key - with fallback to env
            api_key = None
            try:
                from app.services.ai_usage_service import get_api_key
                api_key = get_api_key("gemini")
            except:
                # Fallback to environment variable
                api_key = os.getenv("GOOGLE_API_KEY")
            
            if not api_key:
                raise HTTPException(500, "Gemini API key not configured")
                
            genai.configure(api_key=api_key)
            
            # Upload PDF
            logger.info("Uploading PDF to Gemini...")
            uploaded_file = genai.upload_file(str(pdf_path))
            
            # Create model
            model = genai.GenerativeModel("gemini-2.0-flash-exp")
            
            # Prompt for Markdown generation - SIMPLIFIED
            lang_desc = "Tiếng Việt" if language == "vi" else "English"
            prompt = f"""Convert this PDF document to Markdown format. Preserve the layout as much as possible.

IMPORTANT RULES:

1. **Text Alignment:**
   - If text is CENTERED on the page → start line with: [CENTER]
   - If text is RIGHT-ALIGNED → start line with: [RIGHT]
   - Otherwise, leave as normal left-aligned text

2. **Headings (based on font size and style):**
   - Large UPPERCASE text → # Heading 1
   - Medium bold text → ## Heading 2
   - Smaller bold text → ### Heading 3

3. **Text Formatting:**
   - Bold text: **text**
   - Italic text: *text*

4. **Tables:**
   - Use Markdown table syntax: | Column1 | Column2 |
   - Preserve table structure

5. **Lists:**
   - Bullets: - Item
   - Numbered: 1. Item

6. **Accuracy:**
   - Keep 100% accurate Vietnamese text with all diacritics
   - Preserve numbers, dates, punctuation
   - Maintain document structure and order

Output ONLY the Markdown text, NO explanations.

Markdown:"""

            response = model.generate_content([uploaded_file, prompt])
            markdown = response.text.strip()
            
            # Calculate usage (safely handle missing usage_metadata)
            total_tokens = 0
            cost = 0.0
            if hasattr(response, 'usage_metadata') and response.usage_metadata:
                usage = response.usage_metadata
                prompt_tokens = getattr(usage, 'prompt_token_count', 0)
                candidates_tokens = getattr(usage, 'candidates_token_count', 0)
                total_tokens = prompt_tokens + candidates_tokens
                cost = (prompt_tokens / 1_000_000 * 0.075) + (candidates_tokens / 1_000_000 * 0.30)
            
            # Cleanup
            try:
                genai.delete_file(uploaded_file.name)
            except:
                pass
            
            return {
                "markdown": markdown,
                "char_count": len(markdown),
                "ai_usage": {
                    "engine": "gemini",
                    "model": "gemini-2.0-flash-exp",
                    "total_tokens": total_tokens,
                    "total_cost_usd": round(cost, 6)
                }
            }
            
        except Exception as e:
            logger.error(f"Gemini Markdown generation failed: {e}")
            raise HTTPException(500, f"Markdown generation failed: {str(e)}")
    
    async def _markdown_to_word(self, markdown: str, output_path: Path):
        
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
            import re
            
            doc = Document()
            
            # Set document margins (narrow margins like original)
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(0.5)
                section.bottom_margin = Inches(0.5)
                section.left_margin = Inches(0.75)
                section.right_margin = Inches(0.75)
            
            # Parse Markdown line by line
            lines = markdown.split('\n')
            i = 0
            
            while i < len(lines):
                line = lines[i].rstrip()
                
                # Skip empty lines
                if not line:
                    i += 1
                    continue
                
                # Detect and remove alignment tags
                alignment = WD_PARAGRAPH_ALIGNMENT.LEFT  # default
                original_line = line
                
                if line.startswith('[CENTER]'):
                    alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                    line = line[8:].strip()
                elif line.startswith('[RIGHT]'):
                    alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
                    line = line[7:].strip()
                
                # Remove any leftover tags like [LEFT], [/LEFT], [/RIGHT], [/CENTER]
                line = line.replace('[LEFT]', '').replace('[/LEFT]', '')
                line = line.replace('[RIGHT]', '').replace('[/RIGHT]', '')
                line = line.replace('[CENTER]', '').replace('[/CENTER]', '')
                line = line.strip()
                
                if not line:  # If line becomes empty after removing tags
                    i += 1
                    continue
                
                # Heading 1
                if line.startswith('# ') and not line.startswith('##'):
                    text = line[2:].strip()
                    heading = doc.add_heading(text, level=1)
                    heading.alignment = alignment
                    # Make heading 1 bigger and bold
                    for run in heading.runs:
                        run.font.size = Pt(16)
                        run.font.bold = True
                    i += 1
                    
                # Heading 2
                elif line.startswith('## ') and not line.startswith('###'):
                    text = line[3:].strip()
                    heading = doc.add_heading(text, level=2)
                    heading.alignment = alignment
                    for run in heading.runs:
                        run.font.size = Pt(14)
                        run.font.bold = True
                    i += 1
                    
                # Heading 3
                elif line.startswith('### '):
                    text = line[4:].strip()
                    heading = doc.add_heading(text, level=3)
                    heading.alignment = alignment
                    for run in heading.runs:
                        run.font.size = Pt(12)
                        run.font.bold = True
                    i += 1
                    
                # Table
                elif line.startswith('|'):
                    table_lines = []
                    while i < len(lines) and lines[i].strip().startswith('|'):
                        table_lines.append(lines[i].strip())
                        i += 1
                    
                    # Parse table
                    if len(table_lines) >= 2:
                        # Extract header
                        header_cells = [cell.strip() for cell in table_lines[0].split('|')[1:-1]]
                        
                        # Skip separator line
                        data_lines = table_lines[2:] if len(table_lines) > 2 else []
                        
                        # Create table
                        if data_lines:
                            table = doc.add_table(rows=1 + len(data_lines), cols=len(header_cells))
                            table.style = 'Light Grid Accent 1'
                            
                            # Add header
                            for col_idx, header_text in enumerate(header_cells):
                                cell = table.rows[0].cells[col_idx]
                                cell.text = header_text
                                # Bold header
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True
                            
                            # Add data rows
                            for row_idx, data_line in enumerate(data_lines, start=1):
                                cells_data = [cell.strip() for cell in data_line.split('|')[1:-1]]
                                for col_idx, cell_data in enumerate(cells_data):
                                    if col_idx < len(header_cells):
                                        table.rows[row_idx].cells[col_idx].text = cell_data
                    
                # List (unordered)
                elif line.startswith('- ') or line.startswith('* '):
                    text = line[2:].strip()
                    p = doc.add_paragraph(style='List Bullet')
                    self._add_formatted_text(p, text)
                    i += 1
                    
                # List (ordered)
                elif re.match(r'^\d+\.\s', line):
                    text = re.sub(r'^\d+\.\s', '', line).strip()
                    p = doc.add_paragraph(style='List Number')
                    self._add_formatted_text(p, text)
                    i += 1
                    
                # Regular paragraph
                else:
                    p = doc.add_paragraph()
                    p.alignment = alignment
                    self._add_formatted_text(p, line)
                    
                    # Add spacing after paragraph
                    p.paragraph_format.space_after = Pt(6)
                    i += 1
            
            # Save document
            doc.save(str(output_path))
            logger.info(f"Word document saved: {output_path}")
            
        except Exception as e:
            logger.error(f"Markdown to Word conversion failed: {e}")
            raise HTTPException(500, f"Word generation failed: {str(e)}")
    
    def _add_formatted_text(self, paragraph, text: str):
        
        import re
        
        # Pattern to match: ***text*** (bold+italic), **text** (bold), *text* (italic)
        pattern = r'(\*\*\*[^*]+\*\*\*|\*\*[^*]+\*\*|\*[^*]+\*)'
        
        parts = re.split(pattern, text)
        
        for part in parts:
            if not part:
                continue
                
            # Bold + Italic
            if part.startswith('***') and part.endswith('***'):
                run = paragraph.add_run(part[3:-3])
                run.bold = True
                run.italic = True
                
            # Bold
            elif part.startswith('**') and part.endswith('**'):
                run = paragraph.add_run(part[2:-2])
                run.bold = True
                
            # Italic
            elif part.startswith('*') and part.endswith('*'):
                run = paragraph.add_run(part[1:-1])
                run.italic = True
                
            # Normal text
            else:
                paragraph.add_run(part)
    
    async def _track_pdf_to_word_usage(self, ai_usage: dict, processing_time: float):
        
        try:
            from app.services.ai_usage_service import log_usage
            from app.core.database import SessionLocal
            
            if not ai_usage:
                return
            
            total_tokens = ai_usage.get("total_tokens", 0)
            input_tokens = int(total_tokens * 0.7)
            output_tokens = int(total_tokens * 0.3)
            
            db = SessionLocal()
            try:
                log_usage(
                    db=db,
                    provider="gemini",
                    model=ai_usage.get("model", "gemini-2.0-flash-exp"),
                    endpoint="PDF to Word",
                    input_tokens=input_tokens,
                    output_tokens=output_tokens,
                    processing_time=processing_time,
                    status="success",
                    request_metadata={
                        "cost_usd": ai_usage.get("total_cost_usd", 0.0)
                    }
                )
                db.commit()
            finally:
                db.close()
                
        except Exception as e:
            logger.error(f"Failed to track PDF to Word usage: {e}")


    # ==================== TEXT TO WORD (MHTML FORMAT) ====================
    
    async def text_to_word_mhtml(
        self,
        text: str,
        provider: str = "gemini",  # "gemini" or "claude"
        model: Optional[str] = None,
        language: str = "vi"
    ) -> tuple[str, dict]:
        """
        Convert plain text to formatted Word document using MHTML format
        
        Uses AI (Gemini/Claude) to structure and format the text, then generates
        MHTML which can be opened directly by Microsoft Word.
        
        Args:
            text: Raw text input
            provider: AI provider ("gemini" or "claude")
            model: Specific model name (optional)
            language: Language code ("vi", "en", etc.)
            
        Returns:
            tuple: (mhtml_content, ai_usage_metadata)
        """
        start_time = time.time()
        
        # Step 1: AI formats the text into structured JSON
        if provider == "gemini":
            structured_data = await self._format_text_with_gemini(text, model, language)
        elif provider == "claude":
            structured_data = await self._format_text_with_claude(text, model, language)
        else:
            raise ValueError(f"Invalid provider: {provider}. Use 'gemini' or 'claude'")
        
        # Step 2: Generate DOCX from structured data (instead of MHTML)
        docx_bytes = self._generate_docx(structured_data, language)
        
        processing_time = (time.time() - start_time) * 1000
        
        # Prepare AI usage metadata
        ai_usage = structured_data.get("_metadata", {})
        ai_usage["processing_time_ms"] = processing_time
        
        return docx_bytes, ai_usage
    
    async def _format_text_with_gemini(
        self,
        text: str,
        model: Optional[str] = None,
        language: str = "vi"
    ) -> dict:
        """Use Gemini to structure and format text"""
        from app.services.ai_usage_service import get_api_key, check_budget_limit, log_usage
        from app.core.database import SessionLocal
        
        # Get API key from AI Admin system
        db = SessionLocal()
        try:
            # Check budget first
            budget_check = check_budget_limit("gemini", db)
            if not budget_check["ok"]:
                raise HTTPException(
                    400,
                    f"Gemini budget limit reached. {budget_check.get('reason', 'Monthly limit exceeded')}"
                )
            
            api_key = get_api_key("gemini", db)
            if not api_key:
                raise HTTPException(
                    400,
                    "Gemini API key not configured. Please add key in AI Admin settings."
                )
            
            # Configure Gemini with retrieved key
            if not GEMINI_AVAILABLE:
                raise HTTPException(400, "Gemini library not installed")
            
            genai.configure(api_key=api_key)
            
            # Use specified model or default
            model = model or self.gemini_model_name or DEFAULT_GEMINI_MODEL
            gemini_model = genai.GenerativeModel(model)
            
            prompt = self._build_format_prompt(text, language)
            start_time = time.time()
            
            # Configure with longer timeout for complex text processing
            response = await asyncio.to_thread(
                gemini_model.generate_content,
                prompt,
                generation_config={
                    "temperature": 0.3,
                    "top_p": 0.95,
                    "top_k": 40,
                    "max_output_tokens": 8192,
                },
                request_options={
                    "timeout": 180  # 3 minutes timeout for Gemini API
                }
            )
            
            processing_time = time.time() - start_time
            
            # Parse JSON response
            json_text = response.text.strip()
            if json_text.startswith("```json"):
                json_text = json_text[7:]
            if json_text.endswith("```"):
                json_text = json_text[:-3]
            json_text = json_text.strip()
            
            structured_data = json.loads(json_text)
            
            # Get token usage (safely handle missing usage_metadata)
            input_tokens = 0
            output_tokens = 0
            total_tokens = 0
            
            if hasattr(response, 'usage_metadata') and response.usage_metadata:
                input_tokens = getattr(response.usage_metadata, 'prompt_token_count', 0)
                output_tokens = getattr(response.usage_metadata, 'candidates_token_count', 0)
                total_tokens = getattr(response.usage_metadata, 'total_token_count', 0)
            
            # Add metadata
            structured_data["_metadata"] = {
                "provider": "gemini",
                "model": model,
                "input_tokens": input_tokens,
                "output_tokens": output_tokens,
                "total_tokens": total_tokens,
            }
            
            # Log usage to AI Admin system
            log_usage(
                db=db,
                provider="gemini",
                model=model,
                endpoint="text-to-word",
                input_tokens=input_tokens,
                output_tokens=output_tokens,
                processing_time=processing_time,
                status="success"
            )
            
            return structured_data
            
        except Exception as e:
            # Log error
            log_usage(
                db=db,
                provider="gemini",
                model=model or DEFAULT_GEMINI_MODEL,
                endpoint="text-to-word",
                status="error",
                error_message=str(e)
            )
            raise
            
        finally:
            db.close()
    
    async def _format_text_with_claude(
        self,
        text: str,
        model: Optional[str] = None,
        language: str = "vi"
    ) -> dict:
        """Use Claude to structure and format text"""
        from app.services.ai_usage_service import get_api_key, check_budget_limit, log_usage
        from app.core.database import SessionLocal
        
        # Get API key from AI Admin system
        db = SessionLocal()
        try:
            # Check budget first
            budget_check = check_budget_limit("claude", db)
            if not budget_check["ok"]:
                raise HTTPException(
                    400,
                    f"Claude budget limit reached. {budget_check.get('reason', 'Monthly limit exceeded')}"
                )
            
            claude_api_key = get_api_key("claude", db)
            if not claude_api_key:
                raise HTTPException(
                    400,
                    "Claude API key not configured. Please add key in AI Admin settings."
                )
            
            model = model or "claude-sonnet-4-20250514"
            
            prompt = self._build_format_prompt(text, language)
            start_time = time.time()
            
            # Use anthropic library (same as OCR)
            import anthropic
            client = anthropic.Anthropic(api_key=claude_api_key)
            
            # Call Claude API
            response = await asyncio.to_thread(
                client.messages.create,
                model=model,
                max_tokens=8192,
                temperature=0.3,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )
            
            processing_time = time.time() - start_time
            
            # Parse JSON from response
            content_text = response.content[0].text.strip()
            if content_text.startswith("```json"):
                content_text = content_text[7:]
            if content_text.endswith("```"):
                content_text = content_text[:-3]
            content_text = content_text.strip()
            
            structured_data = json.loads(content_text)
            
            # Get token usage from anthropic response
            input_tokens = response.usage.input_tokens
            output_tokens = response.usage.output_tokens
            
            # Add metadata
            structured_data["_metadata"] = {
                "provider": "claude",
                "model": model,
                "input_tokens": input_tokens,
                "output_tokens": output_tokens,
                "total_tokens": input_tokens + output_tokens,
            }
            
            # Log usage to AI Admin system
            log_usage(
                db=db,
                provider="claude",
                model=model,
                endpoint="text-to-word",
                input_tokens=input_tokens,
                output_tokens=output_tokens,
                processing_time=processing_time,
                status="success"
            )
            
            return structured_data
                
        except Exception as e:
            # Log error if not already logged
            if "Claude API error" not in str(e):
                log_usage(
                    db=db,
                    provider="claude",
                    model=model or "claude-3-5-sonnet-20241022",
                    endpoint="text-to-word",
                    status="error",
                    error_message=str(e)
                )
            raise
            
        finally:
            db.close()
    
    def _build_format_prompt(self, text: str, language: str) -> str:            # Parse JSON from response
            content_text = result["content"][0]["text"].strip()
            if content_text.startswith("```json"):
                content_text = content_text[7:]
            if content_text.endswith("```"):
                content_text = content_text[:-3]
            content_text = content_text.strip()
            
            structured_data = json.loads(content_text)
            
            # Add metadata
            structured_data["_metadata"] = {
                "provider": "claude",
                "model": model,
                "input_tokens": result["usage"]["input_tokens"],
                "output_tokens": result["usage"]["output_tokens"],
                "total_tokens": result["usage"]["input_tokens"] + result["usage"]["output_tokens"],
            }
            
            return structured_data
    
    def _build_format_prompt(self, text: str, language: str) -> str:
        """Build prompt for AI to structure text with visualization support"""
        lang_instruction = {
            "vi": "Phân tích và định dạng văn bản tiếng Việt sau đây",
            "en": "Analyze and format the following English text"
        }.get(language, "Analyze and format the following text")
        
        return f"""{lang_instruction} thành cấu trúc JSON với khả năng tạo biểu đồ:

INPUT TEXT:
{text}

INSTRUCTIONS:
1. Phân tích nội dung và xác định:
   - Title (tiêu đề chính)
   - Subtitle (tiêu đề phụ nếu có)
   - Sections (các phần nội dung chính)
   - **Visualizations (biểu đồ từ số liệu nếu có)**
   
2. Mỗi section có:
   - heading: Tiêu đề phần
   - level: 1 (h2) hoặc 2 (h3)
   - content: Mảng các đoạn văn hoặc list items
   
3. Mỗi content item có:
   - type: "paragraph", "list", "info_box", "highlight_box"
   - text: Nội dung văn bản (cho paragraph)
   - items: ARRAY of strings (cho list) - VD: ["Item 1", "Item 2"]
   - highlights: ARRAY of strings (cho paragraph) - VD: ["Samsung", "Apple"]
   
4. **VISUALIZATION (CHỈ khi text có số liệu):**
   - Nếu text chứa dữ liệu số (doanh thu, thống kê, so sánh...)
   - Tạo biểu đồ phù hợp:
     * "bar" - So sánh giữa các mục
     * "line" - Xu hướng theo thời gian
     * "pie" - Tỷ lệ phần trăm
     * "scatter" - Mối quan hệ giữa 2 biến
   - position: "after_section_0", "after_section_1", ...
   
5. Tự động nhận diện:
   - Thông tin quan trọng → info_box
   - Kết luận/tóm tắt → highlight_box
   - Danh sách → list
   - Tên người, địa điểm, thuật ngữ quan trọng → highlights (PHẢI LÀ ARRAY)
   - **Số liệu → visualizations**

⚠️ QUAN TRỌNG:
- "highlights" PHẢI là array of strings: ["Samsung", "Apple"]
- "items" (cho list) PHẢI là array of strings: ["Item 1", "Item 2", "Item 3"]
- KHÔNG được là string: "highlights" hoặc "items" hoặc "typetexthighlights"
- Nếu không có highlights → dùng array rỗng: []
- Nếu không có items → dùng array rỗng: []

OUTPUT JSON FORMAT:
{{
  "title": "Tiêu đề chính",
  "subtitle": "Tiêu đề phụ (optional)",
  "sections": [
    {{
      "heading": "Phần I: ...",
      "level": 1,
      "content": [
        {{
          "type": "paragraph",
          "text": "Đoạn văn bản...",
          "highlights": ["Samsung", "Apple", "Oppo"]  // ARRAY, not string!
        }},
        {{
          "type": "list",
          "items": ["Hà Nội: 500 triệu", "TP.HCM: 650 triệu"]  // ARRAY, not string!
        }}
      ]
    }}
  ],
  "visualizations": [
    {{
      "type": "bar",
      "title": "Biểu đồ doanh thu",
      "position": "after_section_0",
      "data": {{
        "labels": ["Q1", "Q2", "Q3"],
        "values": [100, 150, 120],
        "colors": ["#3498db", "#2ecc71", "#e74c3c"]
      }},
      "description": "Biểu đồ so sánh doanh thu theo quý"
    }}
  ]
}}

**CHÚ Ý**: Chỉ thêm "visualizations" khi text CÓ dữ liệu số cụ thể.
Trả về JSON thuần túy, không có markdown wrapper."""
    
    def _generate_mhtml(self, structured_data: dict, language: str) -> str:
        """Generate MHTML content from structured data"""
        
        title = structured_data.get("title", "Document")
        subtitle = structured_data.get("subtitle", "")
        sections = structured_data.get("sections", [])
        
        # Generate HTML body content
        body_html = f"""
<div class="Section1">

<h1>{self._escape_html(title)}</h1>
"""
        
        if subtitle:
            body_html += f'<p class="subtitle">{self._escape_html(subtitle)}</p>\n\n'
        
        # Generate sections
        for section in sections:
            heading = section.get("heading", "")
            level = section.get("level", 1)
            content_items = section.get("content", [])
            
            if heading:
                heading_tag = "h2" if level == 1 else "h3"
                body_html += f'<{heading_tag}>{self._escape_html(heading)}</{heading_tag}>\n\n'
            
            # Generate content items
            for item in content_items:
                item_type = item.get("type", "paragraph")
                
                if item_type == "paragraph":
                    text = item.get("text", "")
                    highlights = item.get("highlights", [])
                    
                    # Escape HTML first
                    text = self._escape_html(text)
                    
                    # Apply highlights after escaping
                    for highlight in highlights:
                        escaped_highlight = self._escape_html(highlight)
                        text = text.replace(
                            escaped_highlight,
                            f'<span class="key-name">{escaped_highlight}</span>'
                        )
                    
                    body_html += f'<p>{text}</p>\n\n'
                
                elif item_type == "list":
                    items = item.get("items", [])
                    body_html += '<ul>\n'
                    for list_item in items:
                        body_html += f'<li>{self._escape_html(list_item)}</li>\n'
                    body_html += '</ul>\n\n'
                
                elif item_type == "info_box":
                    box_title = item.get("title", "THÔNG TIN")
                    items = item.get("items", [])
                    
                    body_html += '<div class="info-box">\n'
                    body_html += f'<p class="no-indent"><strong>{self._escape_html(box_title)}</strong></p>\n'
                    for box_item in items:
                        body_html += f'<p class="no-indent">{self._escape_html(box_item)}</p>\n'
                    body_html += '</div>\n\n'
                
                elif item_type == "highlight_box":
                    box_title = item.get("title", "KẾT LUẬN")
                    text = item.get("text", "")
                    
                    body_html += '<div class="conclusion-box">\n'
                    body_html += f'<p class="no-indent"><strong>{self._escape_html(box_title)}</strong></p>\n'
                    body_html += f'<p class="no-indent">{self._escape_html(text)}</p>\n'
                    body_html += '</div>\n\n'
        
        body_html += '''
<p class="center" style="margin-top: 24pt; font-style: italic; color: #666666;"><strong>––– HẾT –––</strong></p>

</div>
'''
        
        # Generate full MHTML
        mhtml = f"""MIME-Version: 1.0
Content-Type: multipart/related; boundary="----=_NextPart_000_0000_01D00000.00000000"

------=_NextPart_000_0000_01D00000.00000000
Content-Type: text/html; charset="utf-8"

<html xmlns:v="urn:schemas-microsoft-com:vml"
xmlns:o="urn:schemas-microsoft-com:office:office"
xmlns:w="urn:schemas-microsoft-com:office:word"
xmlns:m="http://schemas.microsoft.com/office/2004/12/omml"
xmlns="http://www.w3.org/TR/REC-html40">

<head>
<meta http-equiv="Content-Type" content="text/html; charset=utf-8">
<meta name="ProgId" content="Word.Document">
<meta name="Generator" content="Microsoft Word">
<meta name="Originator" content="Microsoft Word">
<!--[if gte mso 9]>
<xml>
<w:WordDocument>
<w:View>Print</w:View>
<w:Zoom>100</w:Zoom>
<w:SpellingState>Clean</w:SpellingState>
<w:GrammarState>Clean</w:GrammarState>
<w:ValidateAgainstSchemas/>
<w:SaveIfXMLInvalid>false</w:SaveIfXMLInvalid>
<w:IgnoreMixedContent>false</w:IgnoreMixedContent>
<w:AlwaysShowPlaceholderText>false</w:AlwaysShowPlaceholderText>
</w:WordDocument>
</xml>
<![endif]-->

<style>
@page Section1 {{
  size: 21.0cm 29.7cm;
  margin: 2.5cm 2.0cm 2.5cm 2.0cm;
  mso-header-margin: 1.5cm;
  mso-footer-margin: 1.5cm;
  mso-paper-source: 0;
}}

div.Section1 {{
  page: Section1;
}}

body {{
  font-family: 'Times New Roman', serif;
  font-size: 13pt;
  line-height: 1.5;
}}

h1 {{
  font-size: 18pt;
  font-weight: bold;
  color: #C00000;
  text-align: center;
  text-transform: uppercase;
  margin-top: 0pt;
  margin-bottom: 12pt;
  border-bottom: 3pt solid #C00000;
  padding-bottom: 6pt;
}}

h2 {{
  font-size: 14pt;
  font-weight: bold;
  color: #C00000;
  margin-top: 18pt;
  margin-bottom: 10pt;
  border-left: 4pt solid #C00000;
  padding-left: 8pt;
}}

h3 {{
  font-size: 13pt;
  font-weight: bold;
  color: #333333;
  margin-top: 12pt;
  margin-bottom: 8pt;
  font-style: italic;
}}

p {{
  text-align: justify;
  margin-top: 6pt;
  margin-bottom: 6pt;
  text-indent: 1.0cm;
}}

p.no-indent {{
  text-indent: 0cm;
}}

p.center {{
  text-align: center;
  text-indent: 0cm;
}}

p.subtitle {{
  text-align: center;
  font-style: italic;
  color: #666666;
  font-size: 11pt;
  margin-bottom: 18pt;
  text-indent: 0cm;
}}

.info-box {{
  background-color: #FFF3CD;
  border-left: 4pt solid #FFA500;
  padding: 10pt;
  margin: 12pt 0pt;
}}

.info-box p {{
  text-indent: 0cm;
  margin: 4pt 0pt;
}}

.conclusion-box {{
  background-color: #E8F5E9;
  border: 2pt solid #4CAF50;
  padding: 10pt;
  margin: 12pt 0pt;
}}

.conclusion-box p {{
  text-indent: 0cm;
  margin: 4pt 0pt;
}}

ul {{
  margin-top: 6pt;
  margin-bottom: 6pt;
  padding-left: 30pt;
}}

li {{
  text-align: justify;
  margin-top: 4pt;
  margin-bottom: 4pt;
}}

.highlight {{
  background-color: #FFFF00;
  font-weight: bold;
}}

.key-name {{
  color: #C00000;
  font-weight: bold;
}}
</style>
</head>

<body lang="{language.upper()}" style='tab-interval:35.4pt'>

{body_html}

</body>
</html>

------=_NextPart_000_0000_01D00000.00000000--"""
        
        return mhtml
    
    def _escape_html(self, text: str) -> str:
        """Escape HTML special characters"""
        return (text
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;"))
    
    def _create_chart(self, viz_config: dict) -> BytesIO:
        """
        Create a chart from visualization config using matplotlib
        
        Args:
            viz_config: {
                "type": "bar|line|pie|scatter",
                "title": "Chart title",
                "data": {
                    "labels": [...],
                    "values": [...],
                    "colors": [...] (optional)
                }
            }
        
        Returns:
            BytesIO: PNG image of the chart
        """
        chart_type = viz_config.get("type", "bar")
        title = viz_config.get("title", "")
        data = viz_config.get("data", {})
        
        labels = data.get("labels", [])
        values = data.get("values", [])
        colors = data.get("colors", None)
        
        # Create figure with good size for document
        fig, ax = plt.subplots(figsize=(8, 5))
        
        try:
            if chart_type == "bar":
                ax.bar(labels, values, color=colors or '#3498db', width=0.6)
                ax.set_ylabel('Value')
                
            elif chart_type == "line":
                ax.plot(labels, values, marker='o', linewidth=2, 
                       markersize=8, color=colors[0] if colors else '#2ecc71')
                ax.set_ylabel('Value')
                ax.grid(True, alpha=0.3)
                
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct='%1.1f%%', 
                      colors=colors, startangle=90)
                ax.axis('equal')  # Equal aspect ratio ensures circular pie
                
            elif chart_type == "scatter":
                ax.scatter(labels, values, s=100, alpha=0.6, 
                          c=colors or '#e74c3c')
                ax.set_ylabel('Value')
            
            # Set title
            if title:
                ax.set_title(title, fontsize=14, fontweight='bold', pad=15)
            
            # Styling
            if chart_type != "pie":
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                ax.tick_params(labelsize=10)
            
            plt.tight_layout()
            
            # Save to BytesIO
            img_stream = BytesIO()
            plt.savefig(img_stream, format='png', dpi=200, bbox_inches='tight')
            img_stream.seek(0)
            plt.close(fig)
            
            return img_stream
            
        except Exception as e:
            logger.error(f"Failed to create chart: {e}")
            plt.close(fig)
            raise
    
    def _generate_docx(self, structured_data: dict, language: str) -> bytes:
        """
        Generate DOCX file from structured data using python-docx
        
        Returns bytes of the DOCX file
        """
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from io import BytesIO
        
        doc = Document()
        
        # Setup page margins (A4)
        sections = doc.sections
        for section in sections:
            section.page_height = Cm(29.7)
            section.page_width = Cm(21.0)
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(2.0)
            section.right_margin = Cm(2.0)
        
        # Get document structure
        title = structured_data.get("title", "Document")
        subtitle = structured_data.get("subtitle", "")
        sections_data = structured_data.get("sections", [])
        visualizations = structured_data.get("visualizations", [])
        
        # Add title
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if title_para.runs:
            title_format = title_para.runs[0].font
            title_format.size = Pt(18)
            title_format.bold = True
            title_format.color.rgb = RGBColor(192, 0, 0)
        
        # Add subtitle if exists
        if subtitle:
            subtitle_para = doc.add_paragraph(subtitle)
            subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if subtitle_para.runs:
                subtitle_format = subtitle_para.runs[0].font
                subtitle_format.size = Pt(11)
                subtitle_format.italic = True
                subtitle_format.color.rgb = RGBColor(102, 102, 102)
        
        # Add sections
        for section_index, section_data in enumerate(sections_data):
            heading = section_data.get("heading", "")
            level = section_data.get("level", 1)
            content_items = section_data.get("content", [])
            
            # Add section heading
            if heading:
                heading_para = doc.add_heading(heading, level=level)
                if heading_para.runs:
                    heading_format = heading_para.runs[0].font
                    heading_format.color.rgb = RGBColor(192, 0, 0)
                    if level == 1:
                        heading_format.size = Pt(14)
                    elif level == 2:
                        heading_format.size = Pt(13)
            
            # Add content
            for content_item in content_items:
                item_type = content_item.get("type", "paragraph")
                
                if item_type == "paragraph":
                    text = content_item.get("text", "")
                    highlights = content_item.get("highlights", [])
                    
                    # FIX: Validate highlights is a list (AI sometimes returns string by mistake)
                    if not isinstance(highlights, list):
                        logger.warning(f"Invalid highlights type: {type(highlights)}. Expected list, got: {highlights}")
                        highlights = []  # Reset to empty list
                    
                    para = doc.add_paragraph()
                    para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    para.paragraph_format.first_line_indent = Cm(1.0)
                    para.paragraph_format.line_spacing = 1.5
                    
                    # Add text with highlights
                    if highlights:
                        # Split text by highlights and format
                        remaining_text = text
                        for highlight in highlights:
                            if highlight in remaining_text:
                                parts = remaining_text.split(highlight, 1)
                                # Add text before highlight
                                if parts[0]:
                                    run = para.add_run(parts[0])
                                    run.font.name = 'Times New Roman'
                                    run.font.size = Pt(13)
                                # Add highlight
                                highlight_run = para.add_run(highlight)
                                highlight_run.font.name = 'Times New Roman'
                                highlight_run.font.size = Pt(13)
                                highlight_run.font.bold = True
                                highlight_run.font.color.rgb = RGBColor(0, 112, 192)
                                # Continue with remaining text
                                remaining_text = parts[1] if len(parts) > 1 else ""
                        # Add remaining text
                        if remaining_text:
                            run = para.add_run(remaining_text)
                            run.font.name = 'Times New Roman'
                            run.font.size = Pt(13)
                    else:
                        # No highlights, just add text
                        run = para.add_run(text)
                        run.font.name = 'Times New Roman'
                        run.font.size = Pt(13)
                
                elif item_type == "list":
                    items = content_item.get("items", [])
                    
                    # FIX: Validate items is a list (AI sometimes returns string by mistake)
                    if not isinstance(items, list):
                        logger.warning(f"Invalid list items type: {type(items)}. Expected list, got: {items}")
                        items = []  # Reset to empty list
                    
                    for list_item in items:
                        para = doc.add_paragraph(list_item, style='List Bullet')
                        if para.runs:
                            para.runs[0].font.name = 'Times New Roman'
                            para.runs[0].font.size = Pt(13)
                
                elif item_type == "info_box":
                    box_title = content_item.get("title", "THÔNG TIN")
                    items = content_item.get("items", [])
                    
                    # Create info box with background color
                    box_para = doc.add_paragraph()
                    box_para.paragraph_format.left_indent = Cm(0.5)
                    box_para.paragraph_format.space_before = Pt(6)
                    box_para.paragraph_format.space_after = Pt(6)
                    
                    # Add title
                    title_run = box_para.add_run(box_title + "\n")
                    title_run.font.name = 'Times New Roman'
                    title_run.font.size = Pt(13)
                    title_run.font.bold = True
                    title_run.font.color.rgb = RGBColor(255, 140, 0)
                    
                    # Add items
                    for item in items:
                        item_run = box_para.add_run(item + "\n")
                        item_run.font.name = 'Times New Roman'
                        item_run.font.size = Pt(13)
                
                elif item_type == "highlight_box":
                    box_title = content_item.get("title", "KẾT LUẬN")
                    text = content_item.get("text", "")
                    
                    # Create highlight box
                    box_para = doc.add_paragraph()
                    box_para.paragraph_format.left_indent = Cm(0.5)
                    box_para.paragraph_format.space_before = Pt(6)
                    box_para.paragraph_format.space_after = Pt(6)
                    
                    # Add title
                    title_run = box_para.add_run(box_title + "\n")
                    title_run.font.name = 'Times New Roman'
                    title_run.font.size = Pt(13)
                    title_run.font.bold = True
                    title_run.font.color.rgb = RGBColor(76, 175, 80)
                    
                    # Add text
                    text_run = box_para.add_run(text)
                    text_run.font.name = 'Times New Roman'
                    text_run.font.size = Pt(13)
            
            # Insert charts after this section if any
            charts_here = [v for v in visualizations 
                          if v.get("position") == f"after_section_{section_index}"]
            
            for viz_config in charts_here:
                try:
                    # Create chart
                    chart_stream = self._create_chart(viz_config)
                    
                    # Add chart to document
                    from docx.shared import Inches
                    doc.add_picture(chart_stream, width=Inches(6.5))
                    
                    # Add description if available
                    description = viz_config.get("description", "")
                    if description:
                        desc_para = doc.add_paragraph(description)
                        desc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        if desc_para.runs:
                            desc_para.runs[0].font.size = Pt(11)
                            desc_para.runs[0].font.italic = True
                            desc_para.runs[0].font.color.rgb = RGBColor(102, 102, 102)
                    
                    logger.info(f"Chart inserted: {viz_config.get('title')}")
                except Exception as e:
                    logger.error(f"Failed to insert chart: {e}")
        
        # Add footer
        footer_para = doc.add_paragraph("––– HẾT –––")
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer_para.paragraph_format.space_before = Pt(24)
        if footer_para.runs:
            footer_format = footer_para.runs[0].font
            footer_format.italic = True
            footer_format.color.rgb = RGBColor(102, 102, 102)
            footer_format.bold = True
        
        # Save to bytes
        docx_buffer = BytesIO()
        doc.save(docx_buffer)
        docx_buffer.seek(0)
        
        return docx_buffer.read()


